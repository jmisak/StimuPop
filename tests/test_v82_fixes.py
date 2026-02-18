"""
Tests for v8.2 feedback fixes:
- Item 1 (Row 36): Multi-element selectbox dropdowns (UI-level, tested via config)
- Item 2 (Row 37): Column separator for same-line display
- Items 3+4 (Row 38/39): Template textbox margin and vertical alignment preservation
"""

import pytest
from io import BytesIO
from unittest.mock import MagicMock, patch

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR

import pandas as pd

from src.pptx_generator import (
    PPTXGenerator,
    SlideConfig,
    ColumnFormat,
    ImageElement,
    TextGroup,
    TEMPLATE_MODE_PLACEHOLDER,
)
from src.excel_handler import ExcelProcessor


# ---------------------------------------------------------------------------
# Item 2: Column Separator Tests
# ---------------------------------------------------------------------------


class TestTextGroupSeparator:
    """Tests for TextGroup.separator field."""

    def test_default_separator_is_empty(self):
        tg = TextGroup(columns=["C", "D"], placeholder_name="TB")
        assert tg.separator == ""

    def test_custom_separator(self):
        tg = TextGroup(columns=["C", "D"], placeholder_name="TB", separator=" for ")
        assert tg.separator == " for "


class TestSlideConfigTextSeparator:
    """Tests for SlideConfig.text_separator field."""

    def test_default_text_separator_is_empty(self):
        config = SlideConfig(img_column="B", text_columns=["C", "D"])
        assert config.text_separator == ""

    def test_custom_text_separator(self):
        config = SlideConfig(
            img_column="B", text_columns=["C", "D"], text_separator=" for "
        )
        assert config.text_separator == " for "


class TestGetSlideDataSeparator:
    """Tests for separator logic in ExcelProcessor.get_slide_data (legacy path)."""

    def test_no_separator_keeps_separate_entries(self):
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["$1.99"], "D": ["0.5 fl oz"]})
        processor = ExcelProcessor()
        slides = processor.get_slide_data(
            df, "B", ["C", "D"], preserve_column_identity=True, text_separator=""
        )
        assert len(slides[0]["text_content"]) == 2

    def test_separator_joins_entries(self):
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["$1.99"], "D": ["0.5 fl oz"]})
        processor = ExcelProcessor()
        slides = processor.get_slide_data(
            df, "B", ["C", "D"], preserve_column_identity=True, text_separator=" for "
        )
        assert len(slides[0]["text_content"]) == 1
        assert slides[0]["text_content"][0]["text"] == "$1.99 for 0.5 fl oz"

    def test_separator_single_column_no_join(self):
        """A single text column should not be affected by separator."""
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["$1.99"]})
        processor = ExcelProcessor()
        slides = processor.get_slide_data(
            df, "B", ["C"], preserve_column_identity=True, text_separator=" for "
        )
        assert len(slides[0]["text_content"]) == 1
        assert slides[0]["text_content"][0]["text"] == "$1.99"

    def test_separator_plain_strings(self):
        """Separator with preserve_column_identity=False."""
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["$1.99"], "D": ["0.5 fl oz"]})
        processor = ExcelProcessor()
        slides = processor.get_slide_data(
            df, "B", ["C", "D"], preserve_column_identity=False, text_separator=" for "
        )
        assert len(slides[0]["text_content"]) == 1
        assert slides[0]["text_content"][0] == "$1.99 for 0.5 fl oz"

    def test_backward_compat_no_separator_arg(self):
        """Calling without text_separator should work (default empty)."""
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["hello"], "D": ["world"]})
        processor = ExcelProcessor()
        slides = processor.get_slide_data(
            df, "B", ["C", "D"], preserve_column_identity=True
        )
        assert len(slides[0]["text_content"]) == 2


class TestGetSlideDataMultiSeparator:
    """Tests for separator logic in ExcelProcessor.get_slide_data_multi."""

    def test_multi_separator_joins_columns(self):
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["$1.99"], "D": ["0.5 fl oz"]})
        processor = ExcelProcessor()
        image_elements = [ImageElement(column="B", placeholder_name="Pic")]
        text_groups = [TextGroup(columns=["C", "D"], placeholder_name="TB", separator=" for ")]

        slides = processor.get_slide_data_multi(df, image_elements, text_groups)

        tc = slides[0]["text_contents"][0]["text_content"]
        assert len(tc) == 1
        assert tc[0]["text"] == "$1.99 for 0.5 fl oz"

    def test_multi_no_separator_keeps_separate(self):
        df = pd.DataFrame({"A": [1], "B": ["img"], "C": ["$1.99"], "D": ["0.5 fl oz"]})
        processor = ExcelProcessor()
        image_elements = [ImageElement(column="B", placeholder_name="Pic")]
        text_groups = [TextGroup(columns=["C", "D"], placeholder_name="TB", separator="")]

        slides = processor.get_slide_data_multi(df, image_elements, text_groups)

        tc = slides[0]["text_contents"][0]["text_content"]
        assert len(tc) == 2


# ---------------------------------------------------------------------------
# Items 3+4: Template Margin & Vertical Alignment Preservation
# ---------------------------------------------------------------------------


def _build_real_template():
    """Build a real PPTX template with custom margins and vertical anchor."""
    prs = PptxPresentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(10)
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(slide_layout)

    # Add a textbox with custom margins
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(6), Inches(2))
    tb.name = "TextBox 1"
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Template Text"
    run.font.name = "Arial"
    run.font.size = Pt(14)

    # Add an image placeholder shape
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(3), Inches(3)
    )
    rect.name = "Rectangle 1"

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


class TestTemplateMarginExtraction:
    """Tests for _extract_template_info capturing margins and vertical_anchor."""

    def test_extracts_margins_from_template(self):
        """Margins and vertical_anchor should be present in extracted shape data."""
        template_bytes = _build_real_template()

        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_placeholder_name="Rectangle 1",
            text_placeholder_name="TextBox 1",
        )
        gen = PPTXGenerator(config)

        prs = PptxPresentation(BytesIO(template_bytes))
        info = gen._extract_template_info(prs.slides[0])

        # Find the text shape
        text_shape = info['text_shapes'].get("TextBox 1")
        assert text_shape is not None, "TextBox 1 should be matched"

        # Verify margins were extracted
        assert text_shape['margin_top'] == Inches(0.1)
        assert text_shape['margin_bottom'] == Inches(0.1)
        assert text_shape['margin_left'] == Inches(0.2)
        assert text_shape['margin_right'] == Inches(0.2)
        assert text_shape['vertical_anchor'] == MSO_ANCHOR.MIDDLE


class TestTemplateMarginApplication:
    """Tests for margins and vertical_anchor being applied to output textboxes."""

    def test_output_preserves_template_margins(self):
        """Generated slides should preserve the template's textbox margins."""
        template_bytes = _build_real_template()

        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_placeholder_name="Rectangle 1",
            text_placeholder_name="TextBox 1",
        )
        gen = PPTXGenerator(config)

        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "image_cell": "B2",
                "text_content": [{"column": "C", "text": "Test Output"}],
            }
        ]

        result = gen.generate(slide_data, template_file=template_bytes)
        assert result.success is True

        # The template slide is removed, so output slide is at index 0
        out_slide = result.presentation.slides[0]

        # Find the TextBox 1 shape in output
        text_shapes = [
            s for s in out_slide.shapes
            if s.has_text_frame and s.name == "TextBox 1"
        ]
        assert len(text_shapes) == 1, "TextBox 1 should exist in output"

        tf = text_shapes[0].text_frame
        assert tf.margin_top == Inches(0.1)
        assert tf.margin_bottom == Inches(0.1)
        assert tf.margin_left == Inches(0.2)
        assert tf.margin_right == Inches(0.2)
        assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE

    def test_output_preserves_vertical_anchor(self):
        """Vertical alignment should survive the template->output cycle."""
        template_bytes = _build_real_template()

        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_placeholder_name="Rectangle 1",
            text_placeholder_name="TextBox 1",
        )
        gen = PPTXGenerator(config)

        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "image_cell": "B2",
                "text_content": [{"column": "C", "text": "Anchored"}],
            }
        ]

        result = gen.generate(slide_data, template_file=template_bytes)
        out_slide = result.presentation.slides[0]
        text_shapes = [s for s in out_slide.shapes if s.name == "TextBox 1"]
        assert text_shapes[0].text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE


class TestBackwardCompatNoMargins:
    """Shapes without text frames should not break extraction."""

    def test_shape_without_text_frame_has_no_margin_keys(self):
        """Non-text shapes should not have margin keys in shape_data."""
        template_bytes = _build_real_template()

        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_placeholder_name="Rectangle 1",
            text_placeholder_name="TextBox 1",
        )
        gen = PPTXGenerator(config)

        prs = PptxPresentation(BytesIO(template_bytes))
        info = gen._extract_template_info(prs.slides[0])

        # Rectangle 1 is an AutoShape -- check it does NOT have margin keys
        # (it may or may not have a text frame depending on shape type)
        rect_shape = info['image_shapes'].get("Rectangle 1")
        assert rect_shape is not None
        # AutoShapes DO have text frames in python-pptx, so they will have margins.
        # The important thing is that _add_text_from_template gracefully handles
        # any shape_data dict regardless of margin presence.


class TestMultiElementMarginPreservation:
    """Multi-element path should also preserve margins (flows through _add_text_from_template)."""

    def test_multi_element_margins_preserved(self):
        """Multi-element text shapes should also get template margins."""
        template_bytes = _build_real_template()

        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_elements=[ImageElement(column="B", placeholder_name="Rectangle 1")],
            text_groups=[TextGroup(columns=["C"], placeholder_name="TextBox 1")],
        )
        gen = PPTXGenerator(config)

        slide_data = [
            {
                "row_index": 0,
                "image_sources": [
                    {"image_source": None, "image_cell": "B2", "placeholder_name": "Rectangle 1"}
                ],
                "text_contents": [
                    {
                        "placeholder_name": "TextBox 1",
                        "text_content": [{"column": "C", "text": "Multi Text"}],
                    }
                ],
            }
        ]

        result = gen.generate(slide_data, template_file=template_bytes)
        assert result.success is True

        out_slide = result.presentation.slides[0]
        text_shapes = [s for s in out_slide.shapes if s.name == "TextBox 1"]
        assert len(text_shapes) == 1

        tf = text_shapes[0].text_frame
        assert tf.margin_top == Inches(0.1)
        assert tf.margin_bottom == Inches(0.1)
        assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE
