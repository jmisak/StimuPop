"""
Tests for PowerPoint generation.
"""

import pytest
from io import BytesIO
from unittest.mock import MagicMock, patch

from PIL import Image

from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.pptx_generator import (
    PPTXGenerator,
    SlideConfig,
    SlideResult,
    GenerationResult,
    ColumnFormat,
    ImageElement,
    TextGroup,
    TEMPLATE_MODE_PLACEHOLDER,
    create_presentation,
)
from src.image_handler import ImageResult
from pptx.dml.color import RGBColor


def create_test_image_result():
    """Create a valid ImageResult for testing."""
    img = Image.new('RGB', (100, 100), color='blue')
    buffer = BytesIO()
    img.save(buffer, format='PNG')
    buffer.seek(0)

    return ImageResult(
        source="test.png",
        success=True,
        data=buffer,
        width=100,
        height=100,
        format="PNG",
        size_bytes=len(buffer.getvalue())
    )


def create_failed_image_result():
    """Create a failed ImageResult for testing."""
    return ImageResult(
        source="failed.png",
        success=False,
        error="Image not found"
    )


class TestColumnFormat:
    """Tests for ColumnFormat dataclass."""

    def test_default_values(self):
        """Test ColumnFormat with minimal required fields."""
        fmt = ColumnFormat(column="C")
        assert fmt.font_size == 14
        assert fmt.bold is False
        assert fmt.italic is False
        assert fmt.font_name == "Calibri"
        assert fmt.color == "000000"

    def test_custom_values(self):
        """Test ColumnFormat with custom values."""
        fmt = ColumnFormat(
            column="D",
            font_size=24,
            bold=True,
            italic=True,
            font_name="Arial",
            color="FF0000"
        )
        assert fmt.column == "D"
        assert fmt.font_size == 24
        assert fmt.bold is True
        assert fmt.italic is True
        assert fmt.font_name == "Arial"
        assert fmt.color == "FF0000"

    def test_get_rgb_color_black(self):
        """Test RGB color conversion for black."""
        fmt = ColumnFormat(column="C", color="000000")
        rgb = fmt.get_rgb_color()
        assert isinstance(rgb, RGBColor)
        assert rgb == RGBColor(0, 0, 0)

    def test_get_rgb_color_red(self):
        """Test RGB color conversion for red."""
        fmt = ColumnFormat(column="C", color="FF0000")
        rgb = fmt.get_rgb_color()
        assert rgb == RGBColor(255, 0, 0)

    def test_get_rgb_color_custom(self):
        """Test RGB color conversion for custom color."""
        fmt = ColumnFormat(column="C", color="1A2B3C")
        rgb = fmt.get_rgb_color()
        assert rgb == RGBColor(26, 43, 60)


class TestSlideConfig:
    """Tests for SlideConfig dataclass."""

    def test_default_values(self):
        """Test SlideConfig with minimal required fields."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C", "D"]
        )
        assert config.img_width == 5.5
        assert config.font_size == 14
        assert config.orientation == "portrait"
        assert config.column_formats is None

    def test_custom_values(self):
        """Test SlideConfig with custom values."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            img_width=6.0,
            font_size=18,
            orientation="landscape"
        )
        assert config.img_width == 6.0
        assert config.font_size == 18
        assert config.orientation == "landscape"

    def test_with_column_formats(self):
        """Test SlideConfig with column_formats."""
        formats = {
            "C": ColumnFormat(column="C", font_size=20, bold=True),
            "D": ColumnFormat(column="D", font_size=12, italic=True),
        }
        config = SlideConfig(
            img_column="B",
            text_columns=["C", "D"],
            column_formats=formats
        )
        assert config.column_formats is not None
        assert len(config.column_formats) == 2

    def test_get_column_format_with_match(self):
        """Test get_column_format returns matching format."""
        fmt_c = ColumnFormat(column="C", font_size=20, bold=True)
        config = SlideConfig(
            img_column="B",
            text_columns=["C", "D"],
            column_formats={"C": fmt_c}
        )
        result = config.get_column_format("C")
        assert result.font_size == 20
        assert result.bold is True

    def test_get_column_format_fallback(self):
        """Test get_column_format returns default for missing column."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C", "D"],
            font_size=16,
            column_formats={"C": ColumnFormat(column="C", font_size=20)}
        )
        # D is not in column_formats, should fall back to default
        result = config.get_column_format("D")
        assert result.font_size == 16  # Falls back to config.font_size
        assert result.bold is False

    def test_get_column_format_no_formats(self):
        """Test get_column_format when column_formats is None."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            font_size=18
        )
        result = config.get_column_format("C")
        assert result.font_size == 18


class TestPPTXGenerator:
    """Tests for PPTXGenerator class."""

    def test_generate_empty_data(self):
        """Test generating with empty slide data."""
        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        result = generator.generate([])

        assert result.success is False
        assert "no slide data" in result.error.lower()

    def test_generate_text_only(self):
        """Test generating slides with text only."""
        config = SlideConfig(img_column="B", text_columns=["C", "D"])
        generator = PPTXGenerator(config)

        slide_data = [
            {"row_index": 0, "image_source": None, "text_content": ["Title 1", "Desc 1"]},
            {"row_index": 1, "image_source": None, "text_content": ["Title 2", "Desc 2"]},
        ]

        result = generator.generate(slide_data)

        assert result.success is True
        assert result.slides_generated == 2
        assert result.slides_with_images == 0
        assert result.presentation is not None

    @patch.object(PPTXGenerator, '_create_slide')
    def test_generate_calls_create_slide(self, mock_create_slide):
        """Test that generate calls _create_slide for each row."""
        mock_create_slide.return_value = SlideResult(index=0, success=True)

        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        slide_data = [
            {"row_index": 0, "image_source": None, "text_content": ["Title"]},
            {"row_index": 1, "image_source": None, "text_content": ["Title"]},
            {"row_index": 2, "image_source": None, "text_content": ["Title"]},
        ]

        result = generator.generate(slide_data)

        assert mock_create_slide.call_count == 3

    def test_generate_with_progress_callback(self):
        """Test progress callback is called."""
        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        callbacks = []

        def callback(status, current, total):
            callbacks.append((status, current, total))

        slide_data = [
            {"row_index": 0, "image_source": None, "text_content": ["Title"]},
        ]

        generator.generate(slide_data, progress_callback=callback)

        assert len(callbacks) > 0

    def test_generate_portrait_orientation(self):
        """Test portrait orientation dimensions."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            orientation="portrait"
        )
        generator = PPTXGenerator(config)

        slide_data = [
            {"row_index": 0, "image_source": None, "text_content": ["Title"]},
        ]

        result = generator.generate(slide_data)

        # Portrait: width < height
        prs = result.presentation
        assert prs.slide_width < prs.slide_height

    def test_generate_landscape_orientation(self):
        """Test landscape orientation dimensions."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            orientation="landscape"
        )
        generator = PPTXGenerator(config)

        slide_data = [
            {"row_index": 0, "image_source": None, "text_content": ["Title"]},
        ]

        result = generator.generate(slide_data)

        # Landscape: width > height
        prs = result.presentation
        assert prs.slide_width > prs.slide_height


class TestSlideCreation:
    """Tests for individual slide creation."""

    def test_slide_with_text(self):
        """Test creating slide with text content."""
        config = SlideConfig(img_column="B", text_columns=["C", "D"])
        generator = PPTXGenerator(config)

        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "text_content": ["Main Title", "Description text"]
            },
        ]

        result = generator.generate(slide_data)

        assert result.success is True
        prs = result.presentation
        assert len(prs.slides) == 1

        # Check that text was added
        slide = prs.slides[0]
        shapes_with_text = [s for s in slide.shapes if hasattr(s, 'text_frame')]
        assert len(shapes_with_text) > 0

    def test_slide_handles_failed_image_gracefully(self):
        """Test that failed images don't crash generation."""
        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        # Create failed image result
        failed_result = create_failed_image_result()
        embedded_images = {"B2": failed_result}

        slide_data = [
            {
                "row_index": 0,
                "image_cell": "B2",
                "image_source": None,
                "text_content": ["Title despite missing image"]
            },
        ]

        result = generator.generate(slide_data, embedded_images=embedded_images)

        # Should still succeed, just without image
        assert result.success is True
        assert result.slides_generated == 1
        assert result.slides_with_images == 0
        assert result.slides_with_errors == 1

    def test_slide_with_embedded_image(self):
        """Test creating slide with embedded image."""
        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        # Create valid image result
        img_result = create_test_image_result()
        embedded_images = {"B2": img_result}

        slide_data = [
            {
                "row_index": 0,
                "image_cell": "B2",
                "image_source": None,
                "text_content": ["Title with image"]
            },
        ]

        result = generator.generate(slide_data, embedded_images=embedded_images)

        assert result.success is True
        assert result.slides_with_images == 1


class TestGenerationResult:
    """Tests for GenerationResult dataclass."""

    def test_default_slide_results(self):
        """Test that slide_results defaults to empty list."""
        result = GenerationResult(success=True)
        assert result.slide_results == []

    def test_with_slide_results(self):
        """Test with explicit slide results."""
        slides = [
            SlideResult(index=0, success=True, has_image=True),
            SlideResult(index=1, success=True, has_image=False),
        ]
        result = GenerationResult(
            success=True,
            slides_generated=2,
            slide_results=slides
        )
        assert len(result.slide_results) == 2


class TestConvenienceFunction:
    """Tests for create_presentation convenience function."""

    def test_create_presentation(self):
        """Test create_presentation function."""
        slide_data = [
            {"row_index": 0, "image_source": None, "text_content": ["Title"]},
        ]

        result = create_presentation(slide_data)

        assert isinstance(result, GenerationResult)
        assert result.success is True


class TestPerColumnFormatting:
    """Tests for per-column text formatting."""

    def test_generate_with_column_formats(self):
        """Test generating slides with per-column formatting."""
        formats = {
            "C": ColumnFormat(column="C", font_size=24, bold=True, color="FF0000"),
            "D": ColumnFormat(column="D", font_size=12, italic=True, color="0000FF"),
        }
        config = SlideConfig(
            img_column="B",
            text_columns=["C", "D"],
            column_formats=formats
        )
        generator = PPTXGenerator(config)

        # New format with column identity preserved
        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "text_content": [
                    {"column": "C", "text": "Title"},
                    {"column": "D", "text": "Description"}
                ]
            },
        ]

        result = generator.generate(slide_data)

        assert result.success is True
        assert result.slides_generated == 1

    def test_generate_with_dict_text_content(self):
        """Test generating slides with dict-style text content."""
        config = SlideConfig(img_column="B", text_columns=["C", "D"])
        generator = PPTXGenerator(config)

        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "text_content": [
                    {"column": "C", "text": "Main Title"},
                    {"column": "D", "text": "Subtitle"}
                ]
            },
        ]

        result = generator.generate(slide_data)

        assert result.success is True
        prs = result.presentation
        assert len(prs.slides) == 1

    def test_backward_compat_string_text_content(self):
        """Test backward compatibility with string text content."""
        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        # Old format: list of strings
        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "text_content": ["Title", "Subtitle"]
            },
        ]

        result = generator.generate(slide_data)

        assert result.success is True
        assert result.slides_generated == 1

    def test_mixed_text_content_formats(self):
        """Test that generator handles mixed text formats."""
        config = SlideConfig(img_column="B", text_columns=["C"])
        generator = PPTXGenerator(config)

        # One slide with new format, one with old
        slide_data = [
            {
                "row_index": 0,
                "image_source": None,
                "text_content": [{"column": "C", "text": "New Format"}]
            },
            {
                "row_index": 1,
                "image_source": None,
                "text_content": ["Old Format"]
            },
        ]

        result = generator.generate(slide_data)

        assert result.success is True
        assert result.slides_generated == 2


# ---------------------------------------------------------------------------
# Helpers for multi-element template tests
# ---------------------------------------------------------------------------

def _make_shape_data(name, shape_type=MSO_SHAPE_TYPE.TEXT_BOX, paragraphs=None):
    """Build a minimal template shape_data dict for testing."""
    return {
        'name': name,
        'type': shape_type,
        'left': Emu(0),
        'top': Emu(0),
        'width': Emu(914400),   # 1 inch
        'height': Emu(914400),  # 1 inch
        'paragraphs': paragraphs or [],
    }


def _make_template_info(shapes, image_shapes=None, text_shapes=None,
                        image_shape=None, text_shape=None):
    """Build a template_info dict matching _extract_template_info output."""
    return {
        'shapes': shapes,
        'image_shape': image_shape,
        'text_shape': text_shape,
        'image_shapes': image_shapes or {},
        'text_shapes': text_shapes or {},
    }


class TestMultiElementTemplateSlide:
    """Tests for v8.0 multi-element _create_slide_from_template."""

    def _generator_with_elements(self, image_elements=None, text_groups=None):
        """Create a PPTXGenerator configured for placeholder mode with elements."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_elements=image_elements,
            text_groups=text_groups,
        )
        return PPTXGenerator(config)

    # --- Multi-image tests ---

    @patch.object(PPTXGenerator, '_add_image_at_position')
    @patch.object(PPTXGenerator, '_recreate_shape')
    def test_multi_image_places_two_images(self, mock_recreate, mock_add_img):
        """Two image_sources entries should each invoke _add_image_at_position."""
        gen = self._generator_with_elements(
            image_elements=[
                ImageElement(column="B", placeholder_name="Picture 1"),
                ImageElement(column="G", placeholder_name="Picture 2"),
            ]
        )

        img1 = create_test_image_result()
        img2 = create_test_image_result()
        embedded = {"B2": img1, "G2": img2}

        pic1_shape = _make_shape_data("Picture 1", MSO_SHAPE_TYPE.AUTO_SHAPE)
        pic2_shape = _make_shape_data("Picture 2", MSO_SHAPE_TYPE.AUTO_SHAPE)
        bg_shape = _make_shape_data("Background")

        template_info = _make_template_info(
            shapes=[pic1_shape, pic2_shape, bg_shape],
            image_shapes={"Picture 1": pic1_shape, "Picture 2": pic2_shape},
        )

        data = {
            "row_index": 0,
            "image_sources": [
                {"image_source": None, "image_cell": "B2", "placeholder_name": "Picture 1"},
                {"image_source": None, "image_cell": "G2", "placeholder_name": "Picture 2"},
            ],
        }

        from pptx import Presentation as _Prs
        prs = _Prs()
        result = gen._create_slide_from_template(prs, data, embedded, template_info)

        assert result.has_image is True
        assert mock_add_img.call_count == 2
        # Background shape should be recreated
        assert mock_recreate.call_count == 1

    # --- Multi-text tests ---

    @patch.object(PPTXGenerator, '_add_text_from_template')
    @patch.object(PPTXGenerator, '_recreate_shape')
    def test_multi_text_places_two_textboxes(self, mock_recreate, mock_add_txt):
        """Two text_contents entries should each invoke _add_text_from_template."""
        gen = self._generator_with_elements(
            text_groups=[
                TextGroup(columns=["C", "D"], placeholder_name="TextBox 1"),
                TextGroup(columns=["E"], placeholder_name="TextBox 2"),
            ]
        )

        tb1_shape = _make_shape_data("TextBox 1")
        tb2_shape = _make_shape_data("TextBox 2")
        bg_shape = _make_shape_data("Background")

        template_info = _make_template_info(
            shapes=[tb1_shape, tb2_shape, bg_shape],
            text_shapes={"TextBox 1": tb1_shape, "TextBox 2": tb2_shape},
        )

        data = {
            "row_index": 0,
            "text_contents": [
                {
                    "placeholder_name": "TextBox 1",
                    "text_content": [
                        {"column": "C", "text": "Title"},
                        {"column": "D", "text": "Subtitle"},
                    ],
                },
                {
                    "placeholder_name": "TextBox 2",
                    "text_content": [{"column": "E", "text": "Footer"}],
                },
            ],
        }

        from pptx import Presentation as _Prs
        prs = _Prs()
        result = gen._create_slide_from_template(prs, data, {}, template_info)

        assert result.text_added is True
        assert mock_add_txt.call_count == 2
        assert mock_recreate.call_count == 1

    # --- Mixed multi-image + multi-text ---

    @patch.object(PPTXGenerator, '_add_image_at_position')
    @patch.object(PPTXGenerator, '_add_text_from_template')
    @patch.object(PPTXGenerator, '_recreate_shape')
    def test_multi_image_and_text_combined(self, mock_recreate, mock_txt, mock_img):
        """Both image_sources and text_contents in one slide."""
        gen = self._generator_with_elements(
            image_elements=[ImageElement(column="B", placeholder_name="Pic")],
            text_groups=[TextGroup(columns=["C"], placeholder_name="TB")],
        )

        img = create_test_image_result()
        embedded = {"B2": img}

        pic_shape = _make_shape_data("Pic", MSO_SHAPE_TYPE.AUTO_SHAPE)
        tb_shape = _make_shape_data("TB")
        dec_shape = _make_shape_data("Decoration")

        template_info = _make_template_info(
            shapes=[pic_shape, tb_shape, dec_shape],
            image_shapes={"Pic": pic_shape},
            text_shapes={"TB": tb_shape},
        )

        data = {
            "row_index": 0,
            "image_sources": [
                {"image_source": None, "image_cell": "B2", "placeholder_name": "Pic"},
            ],
            "text_contents": [
                {"placeholder_name": "TB", "text_content": [{"column": "C", "text": "Hello"}]},
            ],
        }

        from pptx import Presentation as _Prs
        prs = _Prs()
        result = gen._create_slide_from_template(prs, data, embedded, template_info)

        assert result.has_image is True
        assert result.text_added is True
        assert mock_img.call_count == 1
        assert mock_txt.call_count == 1
        assert mock_recreate.call_count == 1

    # --- Placeholder shown when no image data for a shape ---

    @patch.object(PPTXGenerator, '_add_placeholder_shape')
    @patch.object(PPTXGenerator, '_recreate_shape')
    def test_multi_image_no_data_shows_placeholder(self, mock_recreate, mock_placeholder):
        """Image shape with no matching image_sources entry shows placeholder."""
        gen = self._generator_with_elements(
            image_elements=[ImageElement(column="B", placeholder_name="Pic")],
        )

        pic_shape = _make_shape_data("Pic", MSO_SHAPE_TYPE.AUTO_SHAPE)
        template_info = _make_template_info(
            shapes=[pic_shape],
            image_shapes={"Pic": pic_shape},
        )

        # image_sources is present (multi mode) but has no entry for "Pic"
        data = {
            "row_index": 0,
            "image_sources": [],
        }

        from pptx import Presentation as _Prs
        prs = _Prs()
        result = gen._create_slide_from_template(prs, data, {}, template_info)

        assert mock_placeholder.call_count == 1
        assert result.has_image is False

    # --- Legacy backward compatibility ---

    @patch.object(PPTXGenerator, '_add_image_at_position')
    @patch.object(PPTXGenerator, '_add_text_from_template')
    @patch.object(PPTXGenerator, '_recreate_shape')
    def test_legacy_mode_still_works(self, mock_recreate, mock_txt, mock_img):
        """Without image_sources/text_contents, legacy path is used."""
        config = SlideConfig(
            img_column="B",
            text_columns=["C"],
            template_mode=TEMPLATE_MODE_PLACEHOLDER,
            image_placeholder_name="Rectangle 1",
            text_placeholder_name="TextBox",
        )
        gen = PPTXGenerator(config)

        img = create_test_image_result()
        embedded = {"B2": img}

        img_shape = _make_shape_data("Rectangle 1", MSO_SHAPE_TYPE.AUTO_SHAPE)
        txt_shape = _make_shape_data("TextBox 5")

        template_info = _make_template_info(
            shapes=[img_shape, txt_shape],
            image_shape=img_shape,
            text_shape=txt_shape,
            image_shapes={"Rectangle 1": img_shape},
            text_shapes={"TextBox 5": txt_shape},
        )

        data = {
            "row_index": 0,
            "image_source": None,
            "image_cell": "B2",
            "text_content": [{"column": "C", "text": "Legacy text"}],
        }

        from pptx import Presentation as _Prs
        prs = _Prs()
        result = gen._create_slide_from_template(prs, data, embedded, template_info)

        assert result.has_image is True
        assert result.text_added is True
        # Legacy path: image via legacy branch, text via legacy branch
        assert mock_img.call_count == 1
        assert mock_txt.call_count == 1

    # --- Edge: multi-text with empty text_content ---

    @patch.object(PPTXGenerator, '_add_text_from_template')
    def test_multi_text_missing_entry_uses_empty(self, mock_txt):
        """Text shape with no matching text_contents entry passes empty list."""
        gen = self._generator_with_elements(
            text_groups=[TextGroup(columns=["C"], placeholder_name="TB")],
        )

        tb_shape = _make_shape_data("TB")
        template_info = _make_template_info(
            shapes=[tb_shape],
            text_shapes={"TB": tb_shape},
        )

        data = {
            "row_index": 0,
            "text_contents": [],  # Multi mode active but no entry for TB
        }

        from pptx import Presentation as _Prs
        prs = _Prs()
        gen._create_slide_from_template(prs, data, {}, template_info)

        # _add_text_from_template called with empty list
        mock_txt.assert_called_once()
        _, _, text_content_arg = mock_txt.call_args[0]
        assert text_content_arg == []
