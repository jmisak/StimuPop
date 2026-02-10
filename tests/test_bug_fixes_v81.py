"""
Tests for v8.1 bug fixes:
  - BUG 1: Template mode selectbox dropdown population
  - BUG 2: Blank slide image horizontal alignment (img_left)
  - BUG 3: Text alignment and left margin controls
"""

import pytest
from io import BytesIO
from unittest.mock import MagicMock, patch
from PIL import Image

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from src.pptx_generator import (
    SlideConfig,
    ImageAlignment,
    PPTXGenerator,
    IMG_ALIGN_LEFT,
    IMG_ALIGN_CENTER,
    IMG_ALIGN_RIGHT,
    TEMPLATE_MODE_BLANK,
)
from src.image_handler import ImageResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_image_result(width=100, height=100, color="blue"):
    """Create a valid ImageResult with a real PNG image."""
    img = Image.new("RGB", (width, height), color=color)
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return ImageResult(
        source="test.png",
        success=True,
        data=buf,
        width=width,
        height=height,
        format="PNG",
        size_bytes=len(buf.getvalue()),
    )


def _make_config(**overrides):
    """Create a SlideConfig with sensible defaults, accepting overrides."""
    defaults = dict(
        img_column="B",
        text_columns=["C", "D"],
        img_width=5.5,
        img_height=4.0,
        img_top=0.5,
        img_left=0.5,
        text_top=5.0,
        text_left=0.5,
        text_alignment="center",
        font_size=14,
        orientation="portrait",
        template_mode=TEMPLATE_MODE_BLANK,
    )
    defaults.update(overrides)
    return SlideConfig(**defaults)


# ===========================================================================
# BUG 1 Tests: get_template_shape_names
# ===========================================================================

class TestTemplateShapeNames:
    """Verify that get_template_shape_names returns correct shape lists."""

    def _make_template_file(self, shapes):
        """Build a mock .pptx in memory with given shapes on slide 1.

        shapes: list of dicts with 'name' and optionally 'has_text'.
        """
        from pptx.util import Inches as _Inches
        prs = Presentation()
        prs.slide_width = _Inches(10)
        prs.slide_height = _Inches(7.5)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        for s in shapes:
            if s.get("has_text", False):
                slide.shapes.add_textbox(
                    _Inches(0), _Inches(0), _Inches(2), _Inches(1)
                ).name = s["name"]
            else:
                from pptx.enum.shapes import MSO_SHAPE
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    _Inches(0), _Inches(0), _Inches(3), _Inches(3),
                )
                shape.name = s["name"]
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        # Return a mock with .getvalue() like Streamlit UploadedFile
        mock_file = MagicMock()
        mock_file.getvalue.return_value = buf.getvalue()
        return mock_file

    def test_returns_image_and_text_shapes(self):
        from app import get_template_shape_names

        template = self._make_template_file([
            {"name": "Picture 1"},
            {"name": "TextBox 5", "has_text": True},
            {"name": "Rectangle 1"},
        ])
        result = get_template_shape_names(template)
        assert "Picture 1" in result["image_shapes"]
        assert "Rectangle 1" in result["image_shapes"]
        assert "TextBox 5" in result["text_shapes"]

    def test_returns_empty_on_none(self):
        from app import get_template_shape_names

        result = get_template_shape_names(None)
        assert result["image_shapes"] == []
        assert result["text_shapes"] == []

    def test_all_shapes_populated(self):
        from app import get_template_shape_names

        template = self._make_template_file([
            {"name": "Shape A"},
            {"name": "Shape B", "has_text": True},
        ])
        result = get_template_shape_names(template)
        assert "Shape A" in result["all_shapes"]
        assert "Shape B" in result["all_shapes"]


# ===========================================================================
# BUG 2 Tests: Image horizontal alignment (img_left / box_left)
# ===========================================================================

class TestImageHorizontalAlignment:
    """Verify _add_image calculates box_left based on alignment."""

    def _generate_slide_and_get_picture(self, h_align, img_left=0.5):
        """Generate a blank slide and return the picture shape."""
        alignment = ImageAlignment(vertical="center", horizontal=h_align)
        config = _make_config(
            image_alignment=alignment,
            img_left=img_left,
        )
        generator = PPTXGenerator(config)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(slide_layout)

        img_result = _make_image_result(200, 200)
        generator._add_image(slide, prs, img_result)

        # Find the picture shape
        for shape in slide.shapes:
            if shape.shape_type is not None:
                return shape
        return None

    def test_center_alignment_centers_bounding_box(self):
        pic = self._generate_slide_and_get_picture("center", img_left=1.0)
        assert pic is not None
        # Center: box_left = (10 - 5.5) / 2 = 2.25
        # Image is 200x200 (square), scaled to fit 5.5x4.0 box => 4.0x4.0
        # Image centered within box: left = 2.25 + (5.5 - 4.0) / 2 = 2.25 + 0.75 = 3.0
        left_inches = pic.left.inches
        assert 2.9 <= left_inches <= 3.1, f"Expected ~3.0, got {left_inches}"

    def test_left_alignment_uses_img_left(self):
        pic = self._generate_slide_and_get_picture("left", img_left=1.0)
        assert pic is not None
        left_inches = pic.left.inches
        # box_left = img_left = 1.0, image is left-aligned in box => left should be 1.0
        assert 0.9 <= left_inches <= 1.1, f"Expected ~1.0, got {left_inches}"

    def test_right_alignment_uses_img_left_as_margin(self):
        pic = self._generate_slide_and_get_picture("right", img_left=1.0)
        assert pic is not None
        left_inches = pic.left.inches
        # box_left = 10 - 5.5 - 1.0 = 3.5
        # Image is right-aligned in box: left = box_left + box_width - img_width
        # The image fits within 5.5x4.0 box and is 200x200 (square) => scaled to 4.0x4.0
        # So left = 3.5 + 5.5 - 4.0 = 5.0
        assert left_inches >= 3.0, f"Expected rightward placement, got {left_inches}"


# ===========================================================================
# BUG 3 Tests: Text alignment and left margin
# ===========================================================================

class TestSlideConfigTextAlignment:
    """Verify SlideConfig.get_text_pp_align returns correct PP_ALIGN values."""

    def test_default_is_center(self):
        config = _make_config()
        assert config.get_text_pp_align() == PP_ALIGN.CENTER

    def test_left_alignment(self):
        config = _make_config(text_alignment="left")
        assert config.get_text_pp_align() == PP_ALIGN.LEFT

    def test_right_alignment(self):
        config = _make_config(text_alignment="right")
        assert config.get_text_pp_align() == PP_ALIGN.RIGHT

    def test_unknown_falls_back_to_center(self):
        config = _make_config(text_alignment="justify")
        assert config.get_text_pp_align() == PP_ALIGN.CENTER


class TestTextAutoFlowAlignment:
    """Verify _add_text_auto_flow uses config text_alignment and text_left."""

    def _generate_text_slide(self, text_alignment="center", text_left=0.5):
        config = _make_config(text_alignment=text_alignment, text_left=text_left)
        generator = PPTXGenerator(config)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(slide_layout)

        text_items = [
            {"column": "C", "text": "Hello"},
            {"column": "D", "text": "World"},
        ]
        generator._add_text_auto_flow(slide, prs, text_items)
        return slide

    def test_text_left_margin_applied(self):
        slide = self._generate_text_slide(text_left=2.0)
        # The textbox should start at Inches(2.0)
        for shape in slide.shapes:
            if shape.has_text_frame:
                left_inches = shape.left.inches
                assert 1.9 <= left_inches <= 2.1, f"Expected ~2.0, got {left_inches}"
                # Width should be slide_width - text_left - 0.5 = 10 - 2 - 0.5 = 7.5
                width_inches = shape.width.inches
                assert 7.3 <= width_inches <= 7.7, f"Expected ~7.5, got {width_inches}"
                break

    def test_text_alignment_left(self):
        slide = self._generate_text_slide(text_alignment="left")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:  # Skip empty paragraphs
                        assert para.alignment == PP_ALIGN.LEFT
                break

    def test_text_alignment_right(self):
        slide = self._generate_text_slide(text_alignment="right")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        assert para.alignment == PP_ALIGN.RIGHT
                break

    def test_text_alignment_center_default(self):
        slide = self._generate_text_slide(text_alignment="center")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        assert para.alignment == PP_ALIGN.CENTER
                break


class TestTextFixedAlignment:
    """Verify _add_text_fixed uses config text_alignment."""

    def test_fixed_text_uses_config_alignment(self):
        from src.pptx_generator import ColumnPosition

        config = _make_config(
            text_alignment="right",
            column_positions={"C": ColumnPosition(mode="fixed", top=5.0, left=1.0)},
        )
        generator = PPTXGenerator(config)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(slide_layout)

        item = {"column": "C", "text": "Fixed text"}
        col_pos = ColumnPosition(mode="fixed", top=5.0, left=1.0)
        generator._add_text_fixed(slide, prs, item, col_pos)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        assert para.alignment == PP_ALIGN.RIGHT
                break


# ===========================================================================
# Backward compatibility
# ===========================================================================

class TestBackwardCompatibility:
    """Verify defaults preserve legacy behavior."""

    def test_default_img_left(self):
        config = SlideConfig(img_column="B", text_columns=["C"])
        assert config.img_left == 0.5

    def test_default_text_left(self):
        config = SlideConfig(img_column="B", text_columns=["C"])
        assert config.text_left == 0.5

    def test_default_text_alignment(self):
        config = SlideConfig(img_column="B", text_columns=["C"])
        assert config.text_alignment == "center"
        assert config.get_text_pp_align() == PP_ALIGN.CENTER
