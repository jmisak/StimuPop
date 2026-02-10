"""
StimuPop - Excel to PowerPoint Converter

A production-grade web application that converts Excel data to
PowerPoint presentations with embedded images and formatted text.

Features:
- Configurable image alignment (NEW in v6.0)
- Per-column fixed text positioning (NEW in v6.0)
- Simple/Advanced positioning modes (NEW in v6.0)
- Template-based placeholder population
- Configurable paragraph spacing
- Embedded Excel image extraction
- Local file path image support
- Configurable slide layout
- Progress tracking
- Comprehensive error handling

Version: 8.1.0
"""

import tempfile
import os
from io import BytesIO

import streamlit as st
import pandas as pd

from src import (
    Config,
    get_config,
    ExcelProcessor,
    PPTXGenerator,
    SlideConfig,
    ColumnFormat,
    ImageLoader,
    ExcelValidationError,
    PPTXGenerationError,
    IMG_SIZE_FIT_BOX,
    IMG_SIZE_FIT_WIDTH,
    IMG_SIZE_FIT_HEIGHT,
    IMG_SIZE_STRETCH,
    TEMPLATE_MODE_BLANK,
    TEMPLATE_MODE_PLACEHOLDER,
    # NEW in v6.0 - Configurable positioning
    ImageAlignment,
    ColumnPosition,
    IMG_ALIGN_TOP,
    IMG_ALIGN_CENTER,
    IMG_ALIGN_BOTTOM,
    IMG_ALIGN_LEFT,
    IMG_ALIGN_RIGHT,
    # NEW in v8.0 - Multi-element support
    ImageElement,
    TextGroup,
)
from src.logging_config import setup_logging, request_context, get_logger
from src.excel_handler import parse_column_input

# Initialize logging
setup_logging()
logger = get_logger(__name__)


# Page configuration
st.set_page_config(
    page_title="StimuPop v8.1",
    page_icon="🎯",
    layout="wide"
)


def main():
    """Main application entry point."""
    with request_context() as req_id:
        logger.info(f"Session started: {req_id}")
        render_app()


def render_app():
    """Render the main application UI."""
    st.title("🎯 StimuPop v8.1")
    st.markdown("*Excel to PowerPoint with template support*")
    st.markdown("---")

    # Create two columns for layout
    col1, col2 = st.columns([2, 1])

    with col1:
        st.subheader("📁 Upload Files")
        excel_file, template_file = render_file_uploaders()

    with col2:
        st.subheader("⚙️ Configuration")
        img_column, text_columns, font_size, pictures_only = render_basic_config()

    # Advanced settings
    (img_width, img_height, img_size_mode, img_top, text_top, orientation,
     column_formats, paragraph_spacing, template_mode,
     image_placeholder_name, text_placeholder_name,
     img_v_align, img_h_align, column_positions,
     text_overflow_mode, multi_element_enabled, image_elements_config,
     text_groups_config, img_left, text_left,
     text_alignment) = render_advanced_settings(text_columns, font_size, template_file)

    st.markdown("---")

    # Preview Excel data
    df = render_data_preview(excel_file)

    # Preview template if provided
    if template_file and template_mode == TEMPLATE_MODE_PLACEHOLDER:
        render_template_preview(template_file)

    st.markdown("---")

    # Generate button
    render_generate_section(
        excel_file=excel_file,
        template_file=template_file,
        df=df,
        img_column=img_column,
        text_columns=text_columns,
        font_size=font_size,
        img_width=img_width,
        img_height=img_height,
        img_size_mode=img_size_mode,
        img_top=img_top,
        text_top=text_top,
        orientation=orientation,
        column_formats=column_formats,
        paragraph_spacing=paragraph_spacing,
        template_mode=template_mode,
        image_placeholder_name=image_placeholder_name,
        text_placeholder_name=text_placeholder_name,
        img_v_align=img_v_align,
        img_h_align=img_h_align,
        column_positions=column_positions,
        text_overflow_mode=text_overflow_mode,
        multi_element_enabled=multi_element_enabled,
        image_elements_config=image_elements_config,
        text_groups_config=text_groups_config,
        img_left=img_left,
        text_left=text_left,
        text_alignment=text_alignment,
    )

    # Instructions
    render_instructions()

    # Footer
    render_footer()


def render_file_uploaders():
    """Render file upload widgets."""
    config = get_config()

    excel_file = st.file_uploader(
        "Upload Excel File (.xlsx)",
        type=['xlsx'],
        help=f"Upload your Excel file with embedded images (max {config.app.max_upload_size_mb}MB)"
    )

    template_file = st.file_uploader(
        "Upload PowerPoint Template (.pptx)",
        type=['pptx'],
        help="Upload a .pptx template with placeholders for images and text"
    )

    return excel_file, template_file


def get_template_shape_names(template_file) -> dict:
    """Extract shape names from template for UI dropdowns."""
    if not template_file:
        return {"image_shapes": [], "text_shapes": [], "all_shapes": []}
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(BytesIO(template_file.getvalue()))
        if len(prs.slides) == 0:
            return {"image_shapes": [], "text_shapes": [], "all_shapes": []}
        slide = prs.slides[0]
        all_shapes = []
        image_shapes = []
        text_shapes = []
        for shape in slide.shapes:
            all_shapes.append(shape.name)
            if shape.has_text_frame:
                text_shapes.append(shape.name)
            # Shapes that could hold images (rectangles, pictures, placeholders)
            if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PICTURE,
                                     MSO_SHAPE_TYPE.PLACEHOLDER):
                image_shapes.append(shape.name)
            elif not shape.has_text_frame:
                image_shapes.append(shape.name)  # Non-text shapes could be image targets
        return {"image_shapes": image_shapes, "text_shapes": text_shapes, "all_shapes": all_shapes}
    except Exception as e:
        logger.warning(f"Could not extract template shape names: {e}")
        return {"image_shapes": [], "text_shapes": [], "all_shapes": []}


def render_basic_config():
    """Render basic configuration inputs."""
    img_column = st.text_input(
        "Image Column",
        "B",
        help="Column letter or name containing embedded images or file paths (e.g., 'B' or 'Image')"
    )

    # Pictures Only mode (NEW in v6.2)
    pictures_only = st.checkbox(
        "Pictures Only (no text)",
        value=False,
        help="Generate slides with images only, skip text columns"
    )

    if pictures_only:
        text_columns = ""
        font_size = 14
        st.info("📷 Pictures Only mode: Slides will contain images without text")
    else:
        text_columns = st.text_input(
            "Text Columns (comma-separated)",
            "C,D,E,F",
            help="Column letters or names for text content (e.g., 'C,D,E,F')"
        )

        font_size = st.slider(
            "Font Size (pt)",
            min_value=8,
            max_value=48,
            value=14,
            help="Default font size for text content (used in Blank mode)"
        )

    return img_column, text_columns, font_size, pictures_only


def render_advanced_settings(text_columns_str: str, default_font_size: int, template_file=None):
    """Render advanced settings in an expander."""
    column_formats = None
    # Multi-element defaults (always defined for safe return)
    multi_element_enabled = False
    image_elements_config = None
    text_groups_config = None

    # Define size mode options with display labels
    size_mode_options = {
        "Fit to Box (Recommended)": IMG_SIZE_FIT_BOX,
        "Fit Width Only": IMG_SIZE_FIT_WIDTH,
        "Fit Height Only": IMG_SIZE_FIT_HEIGHT,
        "Stretch to Exact Size": IMG_SIZE_STRETCH,
    }

    # Template mode options
    template_mode_options = {
        "Blank Slides (Original)": TEMPLATE_MODE_BLANK,
        "Template Placeholders (NEW)": TEMPLATE_MODE_PLACEHOLDER,
    }

    with st.expander("🔧 Advanced Settings", expanded=True):
        # Template Mode Section (NEW)
        st.markdown("#### 📋 Template Mode")
        template_mode_label = st.selectbox(
            "Generation Mode",
            options=list(template_mode_options.keys()),
            index=0,
            help="Choose how slides are created"
        )
        template_mode = template_mode_options[template_mode_label]

        # Placeholder names (only shown for template mode)
        image_placeholder_name = "Rectangle 1"
        text_placeholder_name = "TextBox"

        if template_mode == TEMPLATE_MODE_PLACEHOLDER:
            st.info("📋 **Template Mode**: Upload a template with placeholder shapes. The first slide will be used as the template.")

            # Get shape names from uploaded template for dropdown population
            shape_names = get_template_shape_names(template_file) if template_file else {
                "image_shapes": [], "text_shapes": [], "all_shapes": []
            }

            placeholder_col1, placeholder_col2 = st.columns(2)
            with placeholder_col1:
                if shape_names["image_shapes"]:
                    # Default to first shape, or "Rectangle 1" if present
                    img_default_idx = 0
                    if "Rectangle 1" in shape_names["image_shapes"]:
                        img_default_idx = shape_names["image_shapes"].index("Rectangle 1")
                    image_placeholder_name = st.selectbox(
                        "Image Placeholder Name",
                        options=shape_names["image_shapes"],
                        index=img_default_idx,
                        help="Select the template shape where images should be placed"
                    )
                else:
                    image_placeholder_name = st.text_input(
                        "Image Placeholder Name",
                        "Rectangle 1",
                        help="Name of the shape where images should be placed (upload template to see dropdown)"
                    )
            with placeholder_col2:
                if shape_names["text_shapes"]:
                    txt_default_idx = 0
                    if "TextBox" in shape_names["text_shapes"]:
                        txt_default_idx = shape_names["text_shapes"].index("TextBox")
                    text_placeholder_name = st.selectbox(
                        "Text Placeholder Name",
                        options=shape_names["text_shapes"],
                        index=txt_default_idx,
                        help="Select the template text box to populate with data"
                    )
                else:
                    text_placeholder_name = st.text_input(
                        "Text Placeholder Name",
                        "TextBox",
                        help="Name (or partial name) of the text box to populate (upload template to see dropdown)"
                    )

            # Multi-Element Mode (NEW in v8.0)
            st.markdown("---")
            multi_element_enabled = st.checkbox(
                "Enable Multi-Element Mode (multiple images/text boxes per slide)",
                value=False,
                help="Map multiple Excel columns to multiple template shapes on the same slide"
            )

            if multi_element_enabled:
                st.markdown("##### Image Elements")
                st.caption("Each image element maps an Excel column to a template shape (use exact shape names from Template Preview)")

                # Initialize/reset session state for element counts
                if "num_image_elements" not in st.session_state:
                    st.session_state.num_image_elements = 1
                if "num_text_groups" not in st.session_state:
                    st.session_state.num_text_groups = 1
            else:
                # Reset counters when multi-element is disabled
                st.session_state.num_image_elements = 1
                st.session_state.num_text_groups = 1

            if multi_element_enabled:

                col_add, col_remove = st.columns(2)
                with col_add:
                    if st.button("Add Image Element", key="add_img_elem"):
                        st.session_state.num_image_elements += 1
                        st.rerun()
                with col_remove:
                    if st.button("Remove Last", key="rm_img_elem") and st.session_state.num_image_elements > 1:
                        st.session_state.num_image_elements -= 1
                        st.rerun()

                image_elements_config = []
                for i in range(st.session_state.num_image_elements):
                    ie_col1, ie_col2 = st.columns(2)
                    with ie_col1:
                        ie_column = st.text_input(
                            f"Image Column #{i+1}", value="B", key=f"ie_col_{i}"
                        )
                    with ie_col2:
                        ie_placeholder = st.text_input(
                            f"Placeholder Name #{i+1}", value="Rectangle 1", key=f"ie_ph_{i}"
                        )
                    image_elements_config.append({
                        "column": ie_column,
                        "placeholder_name": ie_placeholder,
                    })

                st.markdown("##### Text Groups")
                st.caption("Each text group maps one or more Excel columns to a template text box")

                col_add2, col_remove2 = st.columns(2)
                with col_add2:
                    if st.button("Add Text Group", key="add_txt_grp"):
                        st.session_state.num_text_groups += 1
                        st.rerun()
                with col_remove2:
                    if st.button("Remove Last", key="rm_txt_grp") and st.session_state.num_text_groups > 1:
                        st.session_state.num_text_groups -= 1
                        st.rerun()

                text_groups_config = []
                for i in range(st.session_state.num_text_groups):
                    tg_col1, tg_col2 = st.columns(2)
                    with tg_col1:
                        tg_columns = st.text_input(
                            f"Text Columns #{i+1} (comma-separated)",
                            value="C,D,E,F",
                            key=f"tg_cols_{i}",
                        )
                    with tg_col2:
                        tg_placeholder = st.text_input(
                            f"Placeholder Name #{i+1}",
                            value="TextBox",
                            key=f"tg_ph_{i}",
                        )
                    text_groups_config.append({
                        "columns": tg_columns,
                        "placeholder_name": tg_placeholder,
                    })

        else:
            st.info("📄 **Blank Mode**: Creates new blank slides with images and text positioned automatically.")

        st.markdown("---")

        # Paragraph Spacing (NEW)
        st.markdown("#### 📝 Text Spacing")
        paragraph_spacing = st.slider(
            "Paragraph Spacing (points)",
            min_value=0.0,
            max_value=24.0,
            value=0.0,
            step=1.0,
            help="Space after each text paragraph (0 = no extra spacing between columns)"
        )

        # Text Overflow Handling (NEW in v6.2)
        text_overflow_mode = st.selectbox(
            "Text Overflow",
            options=["Resize shape to fit text", "Shrink text on overflow"],
            index=0,
            help="How to handle text that exceeds the text box size"
        )

        st.markdown("---")

        # Image Sizing Section
        st.markdown("#### 🖼️ Image Sizing")
        st.caption("Control how images are sized on each slide")

        size_col1, size_col2 = st.columns(2)

        with size_col1:
            size_mode_label = st.selectbox(
                "Size Mode",
                options=list(size_mode_options.keys()),
                index=0,
                help="How to handle different image sizes"
            )
            img_size_mode = size_mode_options[size_mode_label]

            img_width = st.slider(
                "Max Width (inches)",
                min_value=0.0,
                max_value=9.0,
                value=5.5,
                step=0.25,
                help="Maximum image width (or exact width depending on mode)"
            )

        with size_col2:
            img_height = st.slider(
                "Max Height (inches)",
                min_value=0.0,
                max_value=7.0,
                value=4.0,
                step=0.25,
                help="Maximum image height (used in Fit to Box and Stretch modes)"
            )

            # Show info about current mode
            if img_size_mode == IMG_SIZE_FIT_BOX:
                st.info("Images will fit within the box while maintaining aspect ratio")
            elif img_size_mode == IMG_SIZE_FIT_WIDTH:
                st.info("Images will have fixed width, height adjusts automatically")
            elif img_size_mode == IMG_SIZE_FIT_HEIGHT:
                st.info("Images will have fixed height, width adjusts automatically")
            else:
                st.warning("Images will be stretched to exact size (may distort)")

        # Layout Position - only shown in Blank mode (values ignored in Template mode)
        # Default values (used when Template mode selected)
        img_top = 0.5
        img_left = 0.5
        text_top = 5.0
        text_left = 0.5
        text_alignment = "center"
        orientation = "portrait"

        if template_mode == TEMPLATE_MODE_BLANK:
            st.markdown("---")
            st.markdown("#### 📍 Layout Position (Blank mode only)")
            adv_col1, adv_col2 = st.columns(2)

            with adv_col1:
                img_top = st.slider(
                    "Image Top Position (inches)",
                    min_value=0.0,
                    max_value=3.0,
                    value=0.5,
                    step=0.25,
                    help="Distance from top of slide to image"
                )

                img_left = st.slider(
                    "Image Left Margin (inches)",
                    min_value=0.0,
                    max_value=5.0,
                    value=0.5,
                    step=0.25,
                    help="Distance from left edge of slide to image area (used with Left/Right alignment)"
                )

            with adv_col2:
                text_top = st.slider(
                    "Text Top Position (inches)",
                    min_value=0.0,
                    max_value=8.0,
                    value=5.0,
                    step=0.5,
                    help="Distance from top of slide to text"
                )

                orientation = st.selectbox(
                    "Slide Orientation",
                    options=["portrait", "landscape"],
                    index=0,
                    help="Portrait (tall) or Landscape (wide) slides"
                )

            # Text positioning controls (Blank mode)
            st.markdown("##### Text Box Controls")
            txt_ctrl_col1, txt_ctrl_col2 = st.columns(2)
            with txt_ctrl_col1:
                text_left = st.slider(
                    "Text Left Margin (inches)",
                    min_value=0.0,
                    max_value=5.0,
                    value=0.5,
                    step=0.25,
                    help="Distance from left edge of slide to text area"
                )
            with txt_ctrl_col2:
                text_align_label = st.selectbox(
                    "Text Alignment",
                    options=["Center", "Left", "Right"],
                    index=0,
                    help="Horizontal text alignment within text boxes"
                )
                text_alignment = text_align_label.lower()

        # Image Alignment (NEW in v6.0)
        st.markdown("---")
        st.markdown("#### 🎯 Image Alignment (Blank mode only)")
        st.caption("Control how images are positioned within their bounding box")

        align_col1, align_col2 = st.columns(2)
        with align_col1:
            img_v_align = st.selectbox(
                "Vertical Alignment",
                options=["Center", "Top", "Bottom"],
                index=0,
                help="Vertical position of image within image area"
            )
        with align_col2:
            img_h_align = st.selectbox(
                "Horizontal Alignment",
                options=["Center", "Left", "Right"],
                index=0,
                help="Horizontal position of image within image area"
            )

        # Map display names to constants
        v_align_map = {"Top": "top", "Center": "center", "Bottom": "bottom"}
        h_align_map = {"Left": "left", "Center": "center", "Right": "right"}
        img_v_align_value = v_align_map[img_v_align]
        img_h_align_value = h_align_map[img_h_align]

        # Advanced Positioning Mode (NEW in v6.0)
        st.markdown("---")
        advanced_mode = st.checkbox(
            "🔧 Enable Advanced Positioning",
            value=False,
            help="Enable per-column fixed positioning for precise layout control"
        )

        column_positions = None
        if advanced_mode:
            st.markdown("#### 📐 Per-Column Positioning")
            st.caption("Set fixed positions for columns E and F (common for variety cards)")

            pos_columns = parse_column_input(text_columns_str)
            if pos_columns:
                column_positions = {}
                for col in pos_columns:
                    with st.expander(f"Column {col} Position", expanded=(col in ['E', 'F'])):
                        pos_mode = st.radio(
                            f"Position Mode for {col}",
                            options=["Auto (flow after previous)", "Fixed position"],
                            key=f"pos_mode_{col}",
                            horizontal=True
                        )
                        if pos_mode == "Fixed position":
                            col_top = st.number_input(
                                f"Top position (inches)",
                                min_value=0.0,
                                max_value=10.0,
                                value=5.0 if col == 'E' else 6.5 if col == 'F' else 5.0,
                                step=0.25,
                                key=f"col_top_{col}"
                            )
                            col_left = st.number_input(
                                f"Left margin (inches)",
                                min_value=0.0,
                                max_value=5.0,
                                value=0.5,
                                step=0.25,
                                key=f"col_left_{col}"
                            )
                            column_positions[col] = {
                                "mode": "fixed",
                                "top": col_top,
                                "left": col_left
                            }

        # Per-column formatting section (Blank mode only)
        if template_mode == TEMPLATE_MODE_BLANK:
            st.markdown("---")
            st.markdown("#### 🎨 Column Formatting (Blank mode)")
            column_formats = render_column_format_config(text_columns_str, default_font_size)

    return (img_width, img_height, img_size_mode, img_top, text_top, orientation,
            column_formats, paragraph_spacing, template_mode,
            image_placeholder_name, text_placeholder_name,
            img_v_align_value, img_h_align_value, column_positions,
            text_overflow_mode, multi_element_enabled, image_elements_config,
            text_groups_config, img_left, text_left, text_alignment)


def render_column_format_config(text_columns_str: str, default_font_size: int) -> dict:
    """Render per-column font formatting controls."""
    columns = parse_column_input(text_columns_str)

    if not columns:
        st.caption("Enter text columns above to configure formatting")
        return None

    column_formats = {}

    # Use tabs for each column
    tabs = st.tabs([f"Column {col}" for col in columns])

    font_options = ["Calibri", "Arial", "Times New Roman", "Verdana", "Georgia", "Tahoma"]

    for tab, col in zip(tabs, columns):
        with tab:
            fmt_col1, fmt_col2 = st.columns(2)

            with fmt_col1:
                font_size = st.slider(
                    "Font Size",
                    min_value=8,
                    max_value=48,
                    value=default_font_size,
                    key=f"size_{col}"
                )

                font_name = st.selectbox(
                    "Font",
                    options=font_options,
                    index=0,
                    key=f"font_{col}"
                )

            with fmt_col2:
                color = st.color_picker(
                    "Color",
                    value="#000000",
                    key=f"color_{col}"
                )

                bold = st.checkbox("Bold", key=f"bold_{col}")
                italic = st.checkbox("Italic", key=f"italic_{col}")

            column_formats[col] = ColumnFormat(
                column=col,
                font_size=font_size,
                bold=bold,
                italic=italic,
                font_name=font_name,
                color=color.lstrip("#")
            )

    return column_formats


@st.cache_data(show_spinner=False)
def load_excel_preview(file_bytes: bytes, filename: str):
    """Load Excel file with caching for preview."""
    processor = ExcelProcessor()
    return processor.read_excel(file_bytes, filename)


def load_excel_with_retry(file_bytes: bytes, filename: str, max_retries: int = 2):
    """Load Excel file with retry logic for intermittent failures."""
    last_error = None
    for attempt in range(max_retries + 1):
        try:
            return load_excel_preview(file_bytes, filename)
        except Exception as e:
            last_error = e
            if attempt < max_retries:
                # Clear cache and retry
                load_excel_preview.clear()
                logger.warning(f"Excel load attempt {attempt + 1} failed: {e}, retrying...")
            else:
                raise last_error


def render_data_preview(excel_file):
    """Render Excel data preview."""
    if not excel_file:
        return None

    st.subheader("📋 Data Preview")

    try:
        df = load_excel_with_retry(excel_file.getvalue(), excel_file.name)
        processor = ExcelProcessor()
        summary = processor.get_summary(df)

        # Display summary
        col1, col2 = st.columns(2)
        with col1:
            st.write(f"**Rows:** {summary['row_count']}")
        with col2:
            st.write(f"**Columns:** {summary['column_count']}")

        # Column mapping info
        letter_mapping = ", ".join(
            f"{letter}={name}"
            for letter, name in zip(summary['column_letters'], summary['columns'])
        )
        st.caption(f"Column mapping: {letter_mapping}")

        # Data preview
        st.dataframe(processor.get_preview(df), use_container_width=True)

        return df

    except ExcelValidationError as e:
        st.error(f"❌ {e.message}")
        logger.warning(f"Excel validation error: {e}")
        return None
    except Exception as e:
        st.error(f"❌ Error reading Excel file: {str(e)}")
        logger.error(f"Excel read error: {e}", exc_info=True)
        return None


def render_template_preview(template_file):
    """Preview template structure."""
    st.subheader("📋 Template Preview")

    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(template_file.getvalue()))

        st.write(f"**Slide dimensions:** {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
        st.write(f"**Slides in template:** {len(prs.slides)}")

        if len(prs.slides) > 0:
            slide = prs.slides[0]
            st.write(f"**Shapes in first slide:** {len(slide.shapes)}")

            # Show shape details
            shape_info = []
            for shape in slide.shapes:
                info = {
                    "Name": shape.name,
                    "Type": str(shape.shape_type).split(".")[-1],
                    "Position": f"{shape.left.inches:.2f}\", {shape.top.inches:.2f}\"",
                    "Size": f"{shape.width.inches:.2f}\" × {shape.height.inches:.2f}\""
                }
                if shape.has_text_frame:
                    info["Has Text"] = "Yes"
                shape_info.append(info)

            if shape_info:
                st.dataframe(pd.DataFrame(shape_info), use_container_width=True)

    except Exception as e:
        st.warning(f"Could not preview template: {e}")


def render_generate_section(
    excel_file,
    template_file,
    df,
    img_column,
    text_columns,
    font_size,
    img_width,
    img_height,
    img_size_mode,
    img_top,
    text_top,
    orientation,
    column_formats,
    paragraph_spacing,
    template_mode,
    image_placeholder_name,
    text_placeholder_name,
    img_v_align,
    img_h_align,
    column_positions,
    text_overflow_mode,
    multi_element_enabled=False,
    image_elements_config=None,
    text_groups_config=None,
    img_left=0.5,
    text_left=0.5,
    text_alignment="center",
):
    """Render the generate button and handle generation."""
    if st.button("🎨 Generate Presentation", type="primary", use_container_width=True):
        if not excel_file:
            st.error("❌ Please upload an Excel file first!")
            return

        if df is None or len(df) == 0:
            st.error("❌ Excel file appears to be empty or invalid!")
            return

        if template_mode == TEMPLATE_MODE_PLACEHOLDER and not template_file:
            st.error("❌ Template mode requires a PowerPoint template file!")
            return

        generate_presentation(
            df=df,
            excel_file=excel_file,
            template_file=template_file,
            img_column=img_column,
            text_columns=text_columns,
            font_size=font_size,
            img_width=img_width,
            img_height=img_height,
            img_size_mode=img_size_mode,
            img_top=img_top,
            text_top=text_top,
            orientation=orientation,
            column_formats=column_formats,
            paragraph_spacing=paragraph_spacing,
            template_mode=template_mode,
            image_placeholder_name=image_placeholder_name,
            text_placeholder_name=text_placeholder_name,
            img_v_align=img_v_align,
            img_h_align=img_h_align,
            column_positions=column_positions,
            text_overflow_mode=text_overflow_mode,
            multi_element_enabled=multi_element_enabled,
            image_elements_config=image_elements_config,
            text_groups_config=text_groups_config,
            img_left=img_left,
            text_left=text_left,
            text_alignment=text_alignment,
        )


def generate_presentation(
    df,
    excel_file,
    template_file,
    img_column,
    text_columns,
    font_size,
    img_width,
    img_height,
    img_size_mode,
    img_top,
    text_top,
    orientation,
    column_formats,
    paragraph_spacing,
    template_mode,
    image_placeholder_name,
    text_placeholder_name,
    img_v_align,
    img_h_align,
    column_positions,
    text_overflow_mode,
    multi_element_enabled=False,
    image_elements_config=None,
    text_groups_config=None,
    img_left=0.5,
    text_left=0.5,
    text_alignment="center",
):
    """Generate the PowerPoint presentation."""
    logger.info(f"Starting presentation generation for {excel_file.name} (mode: {template_mode})")

    try:
        processor = ExcelProcessor()

        # Map text overflow mode to python-pptx constant
        overflow_mode = "shrink" if text_overflow_mode == "Shrink text on overflow" else None

        # Multi-element mode (v8.0)
        is_multi = (template_mode == TEMPLATE_MODE_PLACEHOLDER
                    and multi_element_enabled and image_elements_config)

        if is_multi:
            # Validate placeholder names are not empty
            for ie_conf in image_elements_config:
                if not ie_conf.get("placeholder_name", "").strip():
                    st.error("❌ All image element placeholder names must be non-empty")
                    return
            if text_groups_config:
                for tg_conf in text_groups_config:
                    if not tg_conf.get("placeholder_name", "").strip():
                        st.error("❌ All text group placeholder names must be non-empty")
                        return

            # Build ImageElement and TextGroup objects
            image_elements = []
            for ie_conf in image_elements_config:
                image_elements.append(ImageElement(
                    column=ie_conf["column"],
                    placeholder_name=ie_conf["placeholder_name"].strip(),
                    sizing_mode=img_size_mode,
                ))

            text_groups = []
            if text_groups_config:
                for tg_conf in text_groups_config:
                    cols = parse_column_input(tg_conf["columns"])
                    if cols:
                        text_groups.append(TextGroup(
                            columns=cols,
                            placeholder_name=tg_conf["placeholder_name"].strip(),
                        ))

            # Validate multi-element columns
            try:
                resolved_images, resolved_texts = processor.validate_columns_multi(
                    df, image_elements, text_groups
                )
            except ExcelValidationError as e:
                st.error(f"❌ {e.message}")
                if e.details:
                    st.info(f"ℹ️ {e.details}")
                return

            # Extract multi-element slide data
            slide_data = processor.get_slide_data_multi(
                df, image_elements, text_groups
            )

            # Build config with multi-element fields
            config = SlideConfig(
                img_column=image_elements[0].column,
                text_columns=text_groups[0].columns if text_groups else [],
                img_width=img_width,
                img_height=img_height,
                img_size_mode=img_size_mode,
                font_size=font_size,
                paragraph_spacing=paragraph_spacing,
                template_mode=template_mode,
                image_placeholder_name=image_elements[0].placeholder_name,
                text_placeholder_name=text_groups[0].placeholder_name if text_groups else "TextBox",
                text_overflow_mode=overflow_mode,
                image_elements=image_elements,
                text_groups=text_groups,
            )
        else:
            # Legacy single-element path
            text_cols = parse_column_input(text_columns)

            try:
                resolved_img, resolved_text = processor.validate_columns(
                    df, img_column, text_cols
                )
            except ExcelValidationError as e:
                st.error(f"❌ {e.message}")
                if e.details:
                    st.info(f"ℹ️ {e.details}")
                return

            # Extract slide data with column identity preserved
            slide_data = processor.get_slide_data(df, resolved_img, resolved_text, preserve_column_identity=True)

            # Create image alignment config (NEW in v6.0)
            image_alignment = ImageAlignment(
                vertical=img_v_align,
                horizontal=img_h_align
            )

            # Create column positions config (NEW in v6.0)
            col_positions = None
            if column_positions:
                col_positions = {}
                for col, pos in column_positions.items():
                    col_positions[col] = ColumnPosition(
                        mode=pos.get("mode", "auto"),
                        top=pos.get("top"),
                        left=pos.get("left", 0.5)
                    )

            # Create slide configuration with per-column formats
            config = SlideConfig(
                img_column=resolved_img,
                text_columns=resolved_text,
                img_width=img_width,
                img_height=img_height,
                img_size_mode=img_size_mode,
                img_top=img_top,
                img_left=img_left,
                text_top=text_top,
                font_size=font_size,
                orientation=orientation,
                column_formats=column_formats,
                paragraph_spacing=paragraph_spacing,
                template_mode=template_mode,
                image_placeholder_name=image_placeholder_name,
                text_placeholder_name=text_placeholder_name,
                image_alignment=image_alignment,
                column_positions=col_positions,
                text_overflow_mode=overflow_mode,
                text_left=text_left,
                text_alignment=text_alignment,
            )

        # Progress tracking
        progress_bar = st.progress(0)
        status_text = st.empty()

        # Extract embedded images from Excel
        status_text.text("Extracting embedded images...")
        loader = ImageLoader()
        embedded_images = loader.extract_embedded_images(excel_file.getvalue())

        if embedded_images:
            st.info(f"📷 Found {len(embedded_images)} embedded images in Excel")

        # Get template bytes if provided
        template_bytes = template_file.getvalue() if template_file else None

        def progress_callback(status: str, current: int, total: int):
            if total > 0:
                progress_bar.progress(current / total)
            status_text.text(status)

        # Generate presentation
        generator = PPTXGenerator(config)
        result = generator.generate(
            slide_data,
            embedded_images=embedded_images,
            template_file=template_bytes,
            progress_callback=progress_callback
        )

        progress_bar.empty()
        status_text.empty()

        if not result.success:
            st.error(f"❌ Generation failed: {result.error}")
            logger.error(f"Generation failed: {result.error}")
            return

        # Show results summary
        mode_label = "Template" if template_mode == TEMPLATE_MODE_PLACEHOLDER else "Blank"
        if result.slides_with_errors > 0:
            st.warning(
                f"⚠️ Generated {result.slides_generated} slides ({mode_label} mode), "
                f"but {result.slides_with_errors} had issues (images may be missing)"
            )

            # Show detailed errors in expander
            with st.expander("View slide issues"):
                for sr in result.slide_results:
                    if sr.image_error:
                        st.caption(f"Slide {sr.index + 1}: {sr.image_error}")
        else:
            st.success(
                f"✅ Generated {result.slides_generated} slides ({mode_label} mode) "
                f"with {result.slides_with_images} images"
            )

        # Save to temporary file and provide download
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                result.presentation.save(tmp.name)
                tmp_path = tmp.name

            with open(tmp_path, 'rb') as f:
                pptx_data = f.read()
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)

        # Download button
        output_filename = excel_file.name.rsplit('.', 1)[0] + '_presentation.pptx'
        st.download_button(
            label="📥 Download Presentation",
            data=pptx_data,
            file_name=output_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

        logger.info(f"Presentation generated successfully: {output_filename}")

    except Exception as e:
        st.error(f"❌ Error generating presentation: {str(e)}")
        logger.error(f"Generation error: {e}", exc_info=True)
        st.exception(e)


def render_instructions():
    """Render usage instructions."""
    st.markdown("---")
    with st.expander("📖 How to Use"):
        st.markdown("""
### Instructions

#### Template Mode

1. **Prepare your template:**
   - Create a PowerPoint slide with placeholder shapes
   - Name the image placeholder (e.g., "Rectangle 1")
   - Name the text placeholder (e.g., "TextBox 55")
   - The first slide will be cloned for each row

2. **Configure placeholders:**
   - Enter the exact shape names in Advanced Settings
   - Text will populate preserving the template's formatting

#### Blank Mode (Original)

1. **Prepare your Excel file:**
   - Each row becomes one slide
   - Include images embedded in a column (e.g., column B) or file paths to local images
   - Include columns with text content (e.g., columns C, D, E, F)

2. **Upload your files:**
   - Upload the Excel file with embedded images (required)
   - Optionally upload a PowerPoint template

3. **Configure settings:**
   - Specify which column contains images
   - Specify which columns contain text
   - Adjust font size, spacing, and positioning as needed

4. **Generate:**
   - Click "Generate Presentation"
   - Download your completed PowerPoint file

### Column Reference

You can reference columns by:
- **Letter**: A, B, C, ... (Excel-style)
- **Name**: The actual column header name

### Tips

- **Template mode** preserves your slide design and formatting
- **Paragraph spacing = 0** removes gaps between text lines
- Embed images directly in Excel cells for best results
- Or use local file paths (e.g., `C:\\Images\\photo.jpg`)
- Images are automatically centered and sized
- Missing images are skipped (slide still created)
        """)


def render_footer():
    """Render footer."""
    st.markdown("---")
    st.markdown(
        "<p style='text-align: center; color: gray;'>"
        "🎯 StimuPop v8.1.0 |"
        "Image/Text Alignment Fixes | "
        "Built with Streamlit"
        "</p>",
        unsafe_allow_html=True
    )


if __name__ == "__main__":
    main()
