"""
PowerPoint generation for StimuPop v8.1.

Provides presentation creation with:
- Template-based generation
- Placeholder detection and population
- Configurable paragraph spacing
- Configurable image alignment (v6.0)
- Per-column fixed positioning (v6.0)
- Embedded image support
- Local file path support
- Configurable layout
- Progress callbacks
- Error handling per slide
"""

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple, Union
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from .config import get_config
from .exceptions import PPTXGenerationError
from .image_handler import ImageLoader, ImageResult
from .logging_config import get_logger

logger = get_logger(__name__)


@dataclass
class ColumnFormat:
    """
    Font formatting for a single text column.

    Attributes:
        column: Column identifier (letter or name)
        font_size: Font size in points
        bold: Whether text is bold
        italic: Whether text is italic
        font_name: Font family name
        color: Hex color without # (e.g., "000000" for black)
    """
    column: str
    font_size: int = 14
    bold: bool = False
    italic: bool = False
    font_name: str = "Calibri"
    color: str = "000000"

    def get_rgb_color(self) -> RGBColor:
        """Convert hex color to RGBColor."""
        r = int(self.color[0:2], 16)
        g = int(self.color[2:4], 16)
        b = int(self.color[4:6], 16)
        return RGBColor(r, g, b)


# Image sizing modes
IMG_SIZE_FIT_BOX = "fit_box"      # Fit within width AND height, maintain aspect ratio
IMG_SIZE_FIT_WIDTH = "fit_width"  # Fixed width, auto height (original behavior)
IMG_SIZE_FIT_HEIGHT = "fit_height"  # Fixed height, auto width
IMG_SIZE_STRETCH = "stretch"      # Exact size, may distort

# Template modes
TEMPLATE_MODE_BLANK = "blank"      # Create blank slides (original behavior)
TEMPLATE_MODE_PLACEHOLDER = "placeholder"  # Use template placeholders

# Image alignment options (NEW in v6.0)
IMG_ALIGN_TOP = "top"
IMG_ALIGN_CENTER = "center"
IMG_ALIGN_BOTTOM = "bottom"
IMG_ALIGN_LEFT = "left"
IMG_ALIGN_RIGHT = "right"


@dataclass
class ImageAlignment:
    """
    Controls image positioning within the image bounding box.

    Attributes:
        vertical: 'top' | 'center' | 'bottom' - vertical alignment within box
        horizontal: 'left' | 'center' | 'right' - horizontal alignment within box
    """
    vertical: str = IMG_ALIGN_CENTER
    horizontal: str = IMG_ALIGN_CENTER


@dataclass
class ColumnPosition:
    """
    Per-column text positioning configuration.

    Attributes:
        mode: 'auto' | 'fixed' - auto flows after previous, fixed uses explicit position
        top: Fixed top position in inches (only used if mode='fixed')
        left: Left margin in inches
        width: Text box width in inches (None = slide width - margins)
    """
    mode: str = "auto"  # auto | fixed
    top: Optional[float] = None  # Only used if mode='fixed'
    left: float = 0.5  # Default left margin
    width: Optional[float] = None  # None = auto width


@dataclass
class ImageElement:
    """
    Configuration for a single image element on the slide (NEW in v8.0).

    Attributes:
        column: Excel column letter containing the image (e.g., "B")
        placeholder_name: Template shape name to place the image into (e.g., "Picture 3")
        sizing_mode: How to size the image (fit_box, fit_width, fit_height, stretch)
        alignment: Optional image alignment within the placeholder bounds
    """
    column: str
    placeholder_name: str
    sizing_mode: str = IMG_SIZE_FIT_BOX
    alignment: Optional[ImageAlignment] = None


@dataclass
class TextGroup:
    """
    Configuration for a group of text columns mapped to a single text box (NEW in v8.0).

    Attributes:
        columns: Excel column letters to populate this text box (e.g., ["C", "D"])
        placeholder_name: Template shape name to populate (e.g., "TextBox 5")
    """
    columns: List[str]
    placeholder_name: str


@dataclass
class SlideConfig:
    """
    Configuration for slide generation.

    Attributes:
        img_column: Column containing images (embedded or file paths)
        text_columns: Columns containing text content
        img_width: Image width in inches (max width for fit_box mode)
        img_height: Image height in inches (max height for fit_box mode)
        img_size_mode: How to size images (fit_box, fit_width, fit_height, stretch)
        img_top: Image top position in inches
        text_top: Text top position in inches
        font_size: Font size in points (fallback default)
        orientation: 'portrait' or 'landscape'
        column_formats: Per-column font formatting
        paragraph_spacing: Space after each paragraph in points (NEW)
        template_mode: 'blank' or 'placeholder' (NEW)
        image_placeholder_name: Name of image placeholder shape (NEW)
        text_placeholder_name: Name of text placeholder shape (NEW)
        image_elements: List of multi-image configs (NEW in v8.0)
        text_groups: List of multi-text-box configs (NEW in v8.0)
    """
    img_column: str
    text_columns: List[str]
    img_width: float = 5.5
    img_height: float = 4.0
    img_size_mode: str = IMG_SIZE_FIT_BOX
    img_top: float = 0.5
    text_top: float = 5.0
    font_size: int = 14
    orientation: str = "portrait"
    column_formats: Optional[Dict[str, ColumnFormat]] = None
    # NEW in v5.1
    paragraph_spacing: float = 0.0  # Points - 0 means no extra spacing
    template_mode: str = TEMPLATE_MODE_BLANK
    image_placeholder_name: str = "Rectangle 1"  # Default from Variety Card
    text_placeholder_name: str = "TextBox"  # Will match any TextBox
    # NEW in v6.0 - Configurable positioning
    image_alignment: Optional[ImageAlignment] = None  # None = center (legacy behavior)
    column_positions: Optional[Dict[str, ColumnPosition]] = None  # None = auto flow all
    positioning_mode: str = "simple"  # simple | advanced
    # NEW in v8.1 - Blank mode image left margin
    img_left: float = 0.5  # Left margin for image bounding box in inches (blank mode)
    # NEW in v6.2 - Text overflow handling
    text_overflow_mode: Optional[str] = None  # None = resize shape, "shrink" = shrink text
    # NEW in v8.1 - Text alignment and left margin (Blank mode)
    text_alignment: str = "center"  # "left" | "center" | "right"
    text_left: float = 0.5  # Left margin for text boxes in inches (blank mode)
    # NEW in v8.0 - Multi-element support
    image_elements: Optional[List['ImageElement']] = None
    text_groups: Optional[List['TextGroup']] = None

    def get_column_format(self, column: str) -> ColumnFormat:
        """Get format for column, falling back to defaults."""
        if self.column_formats and column in self.column_formats:
            return self.column_formats[column]
        return ColumnFormat(column=column, font_size=self.font_size)

    def get_image_alignment(self) -> ImageAlignment:
        """Get image alignment, defaulting to center if not set."""
        if self.image_alignment is not None:
            return self.image_alignment
        return ImageAlignment()  # Default: center/center

    def get_text_pp_align(self) -> PP_ALIGN:
        """Get PowerPoint paragraph alignment from text_alignment config."""
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        return align_map.get(self.text_alignment, PP_ALIGN.CENTER)

    def get_column_position(self, column: str) -> Optional[ColumnPosition]:
        """Get position config for column, or None if auto."""
        if self.column_positions and column in self.column_positions:
            return self.column_positions[column]
        return None  # Auto flow

    def get_image_elements(self) -> List['ImageElement']:
        """Get image elements list, building from legacy fields if not set (v8.0 compat)."""
        if self.image_elements is not None:
            return self.image_elements  # Explicit list (even if empty)
        # Backward compat: build single-item list from legacy fields
        return [ImageElement(
            column=self.img_column,
            placeholder_name=self.image_placeholder_name,
            sizing_mode=self.img_size_mode,
            alignment=self.image_alignment,
        )]

    def get_text_groups(self) -> List['TextGroup']:
        """Get text groups list, building from legacy fields if not set (v8.0 compat)."""
        if self.text_groups is not None:
            return self.text_groups  # Explicit list (even if empty)
        # Backward compat: build single-item list from legacy fields
        if self.text_columns:
            return [TextGroup(
                columns=self.text_columns,
                placeholder_name=self.text_placeholder_name,
            )]
        return []


@dataclass
class SlideResult:
    """Result of generating a single slide."""
    index: int
    success: bool
    has_image: bool = False
    image_error: Optional[str] = None
    text_added: bool = False
    error: Optional[str] = None


@dataclass
class GenerationResult:
    """Result of generating a presentation."""
    success: bool
    presentation: Optional[Presentation] = None
    slides_generated: int = 0
    slides_with_images: int = 0
    slides_with_errors: int = 0
    slide_results: List[SlideResult] = None
    error: Optional[str] = None

    def __post_init__(self):
        if self.slide_results is None:
            self.slide_results = []


class PPTXGenerator:
    """
    Generates PowerPoint presentations from slide data.

    Features:
    - Template-based generation with placeholder support (NEW)
    - Configurable paragraph spacing (NEW)
    - Embedded Excel images
    - Local file path images
    - Configurable slide layout
    - Progress callbacks
    - Per-slide error handling (continues on error)
    """

    def __init__(
        self,
        config: Optional[SlideConfig] = None,
        image_loader: Optional[ImageLoader] = None
    ):
        """Initialize generator."""
        app_config = get_config()
        pres_config = app_config.presentation

        if config is None:
            config = SlideConfig(
                img_column="B",
                text_columns=["C", "D", "E", "F"],
                img_width=pres_config.default_img_width,
                img_top=pres_config.default_img_top,
                text_top=pres_config.default_text_top,
                font_size=pres_config.default_font_size,
                orientation=pres_config.default_orientation
            )

        self.config = config
        self.loader = image_loader or ImageLoader()
        self.app_config = app_config
        self._template_shapes = None  # Cache template shape data

    def generate(
        self,
        slide_data: List[dict],
        embedded_images: Optional[Dict[str, ImageResult]] = None,
        template_file: Optional[Union[bytes, BytesIO, str, Path]] = None,
        progress_callback: Optional[Callable[[str, int, int], None]] = None
    ) -> GenerationResult:
        """Generate a PowerPoint presentation."""
        if not slide_data:
            return GenerationResult(
                success=False,
                error="No slide data provided"
            )

        embedded_images = embedded_images or {}

        try:
            # Create or load presentation
            prs, template_info = self._create_presentation(template_file)

            total_slides = len(slide_data)

            if progress_callback:
                progress_callback("Creating slides...", 0, total_slides)

            # Generate slides
            slide_results = []
            slides_with_images = 0
            slides_with_errors = 0

            for i, data in enumerate(slide_data):
                if progress_callback:
                    progress_callback(
                        f"Creating slide {i + 1}/{total_slides}",
                        i + 1,
                        total_slides
                    )

                # Choose generation method based on template mode
                if self.config.template_mode == TEMPLATE_MODE_PLACEHOLDER and template_info:
                    result = self._create_slide_from_template(prs, data, embedded_images, template_info)
                else:
                    result = self._create_slide(prs, data, embedded_images)

                slide_results.append(result)

                if result.has_image:
                    slides_with_images += 1
                if result.image_error or result.error:
                    slides_with_errors += 1

            # Remove the template slide if we used placeholder mode
            if self.config.template_mode == TEMPLATE_MODE_PLACEHOLDER and template_info:
                self._remove_slide(prs, 0)

            logger.info(
                f"Generated presentation: {total_slides} slides, "
                f"{slides_with_images} with images, "
                f"{slides_with_errors} with errors"
            )

            return GenerationResult(
                success=True,
                presentation=prs,
                slides_generated=total_slides,
                slides_with_images=slides_with_images,
                slides_with_errors=slides_with_errors,
                slide_results=slide_results
            )

        except Exception as e:
            logger.error(f"Presentation generation failed: {e}", exc_info=True)
            return GenerationResult(
                success=False,
                error=str(e)
            )

    def _create_presentation(
        self,
        template_file: Optional[Union[bytes, BytesIO, str, Path]]
    ) -> Tuple[Presentation, Optional[dict]]:
        """Create or load a presentation. Returns (presentation, template_info)."""
        pres_config = self.app_config.presentation
        template_info = None

        if template_file:
            try:
                if isinstance(template_file, bytes):
                    template_file = BytesIO(template_file)
                prs = Presentation(template_file)

                # If using placeholder mode, extract template info
                if self.config.template_mode == TEMPLATE_MODE_PLACEHOLDER and len(prs.slides) > 0:
                    template_info = self._extract_template_info(prs.slides[0])
                    logger.info(f"Extracted template info: {len(template_info.get('shapes', []))} shapes")

                logger.info(f"Loaded template: {prs.slide_width.inches:.2f}x{prs.slide_height.inches:.2f} inches")
            except Exception as e:
                raise PPTXGenerationError(
                    f"Cannot load template: {e}",
                    operation="load_template"
                )
        else:
            prs = Presentation()

            # Set dimensions based on orientation
            if self.config.orientation == "landscape":
                prs.slide_width = Inches(pres_config.landscape_width_inches)
                prs.slide_height = Inches(pres_config.landscape_height_inches)
            else:
                prs.slide_width = Inches(pres_config.portrait_width_inches)
                prs.slide_height = Inches(pres_config.portrait_height_inches)

            logger.info(
                f"Created new presentation: "
                f"{prs.slide_width.inches}x{prs.slide_height.inches} inches"
            )

        return prs, template_info

    def _extract_template_info(self, template_slide) -> dict:
        """Extract shape information from template slide.

        Supports both legacy single-element mode and v8.0 multi-element mode.
        In multi-element mode, image_elements and text_groups on config drive
        which shapes are matched. Legacy fields (image_shape, text_shape) are
        kept as aliases to the first match for backward compatibility.
        """
        info = {
            'shapes': [],
            'image_shape': None,       # Legacy: first matched image shape
            'text_shape': None,        # Legacy: first matched text shape
            'image_shapes': {},        # v8.0: {shape_name: shape_data}
            'text_shapes': {},         # v8.0: {shape_name: shape_data}
        }

        # Determine matching mode: exact for multi-element, substring for legacy
        use_exact_match = bool(self.config.image_elements or self.config.text_groups)

        # Build lookup sets for matching
        if self.config.image_elements:
            image_placeholder_names = [
                el.placeholder_name.lower() for el in self.config.image_elements
                if el.placeholder_name.strip()  # Guard: skip empty names
            ]
        else:
            image_placeholder_names = [self.config.image_placeholder_name.lower()]

        if self.config.text_groups:
            text_placeholder_names = [
                tg.placeholder_name.lower() for tg in self.config.text_groups
                if tg.placeholder_name.strip()  # Guard: skip empty names
            ]
        else:
            text_placeholder_names = [self.config.text_placeholder_name.lower()]

        for shape in template_slide.shapes:
            shape_data = {
                'name': shape.name,
                'type': shape.shape_type,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'paragraphs': []
            }

            # Extract text frame info
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_data = {
                        'text': para.text,
                        'alignment': para.alignment,
                        'level': para.level,
                        'runs': []
                    }
                    for run in para.runs:
                        run_data = {
                            'text': run.text,
                            'font_name': run.font.name,
                            'font_size': run.font.size,
                            'bold': run.font.bold,
                            'italic': run.font.italic,
                        }
                        try:
                            run_data['color'] = run.font.color.rgb
                        except (AttributeError, TypeError):
                            run_data['color'] = None
                        para_data['runs'].append(run_data)
                    shape_data['paragraphs'].append(para_data)

            info['shapes'].append(shape_data)

            # Match image shapes
            shape_name_lower = shape.name.lower()
            for placeholder in image_placeholder_names:
                matched = (shape_name_lower == placeholder) if use_exact_match else (placeholder in shape_name_lower)
                if matched:
                    info['image_shapes'][shape.name] = shape_data
                    if info['image_shape'] is None:
                        info['image_shape'] = shape_data
                    break

            # Match text shapes (only if not already matched as image)
            if shape.name not in info['image_shapes']:
                for placeholder in text_placeholder_names:
                    matched = (shape_name_lower == placeholder) if use_exact_match else (placeholder in shape_name_lower)
                    if matched:
                        info['text_shapes'][shape.name] = shape_data
                        if info['text_shape'] is None:
                            info['text_shape'] = shape_data
                        break

        logger.debug(
            f"Template extraction: {len(info['image_shapes'])} image shape(s), "
            f"{len(info['text_shapes'])} text shape(s) matched"
        )

        return info

    def _remove_slide(self, prs: Presentation, slide_idx: int) -> None:
        """Remove a slide by index."""
        try:
            slide_id = prs.slides._sldIdLst[slide_idx].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[slide_idx]
            logger.info(f"Removed slide at index {slide_idx}")
        except Exception as e:
            logger.warning(f"Could not remove slide: {e}")

    def _create_slide_from_template(
        self,
        prs: Presentation,
        data: dict,
        embedded_images: Dict[str, ImageResult],
        template_info: dict
    ) -> SlideResult:
        """Create a slide by recreating template shapes and populating with data.

        Supports both legacy single-element mode and v8.0 multi-element mode.

        Multi-element mode is active when data contains 'image_sources' or
        'text_contents' lists. Each entry carries a 'placeholder_name' that is
        matched against the template shape name.

        Legacy mode uses the single 'image_source'/'image_cell'/'text_content'
        fields and template_info['image_shape']/['text_shape'].
        """
        index = data.get("row_index", len(prs.slides))
        result = SlideResult(index=index, success=True)

        try:
            # Use blank layout to avoid unwanted placeholder shapes
            blank_layout = None
            for layout in prs.slide_layouts:
                if layout.name == "Blank" or "blank" in layout.name.lower():
                    blank_layout = layout
                    break
            if blank_layout is None:
                blank_layout = prs.slide_layouts[-1]  # Use last layout as fallback

            new_slide = prs.slides.add_slide(blank_layout)

            # Detect multi-element mode: presence of key (even empty list) means multi
            is_multi = "image_sources" in data or "text_contents" in data

            if is_multi:
                self._populate_slide_multi(
                    new_slide, prs, data, embedded_images, template_info, result
                )
            else:
                self._populate_slide_legacy(
                    new_slide, prs, data, embedded_images, template_info, result
                )

        except Exception as e:
            result.success = False
            result.error = str(e)
            logger.error(f"Slide {index}: Creation failed: {e}", exc_info=True)

        return result

    # ------------------------------------------------------------------
    # Multi-element slide population (v8.0)
    # ------------------------------------------------------------------

    def _populate_slide_multi(
        self,
        slide,
        prs: Presentation,
        data: dict,
        embedded_images: Dict[str, ImageResult],
        template_info: dict,
        result: SlideResult
    ) -> None:
        """Populate a slide using v8.0 multi-element data.

        Walks every template shape and dispatches to image, text, or
        recreate based on the image_shapes / text_shapes lookup dicts.
        """
        image_shapes = template_info.get('image_shapes', {})
        text_shapes = template_info.get('text_shapes', {})

        # Build fast lookup: lowercase placeholder_name -> image element data
        image_source_map: Dict[str, dict] = {}
        for img_entry in (data.get("image_sources") or []):
            image_source_map[img_entry["placeholder_name"].lower()] = img_entry

        # Build fast lookup: lowercase placeholder_name -> text group data
        text_content_map: Dict[str, dict] = {}
        for txt_entry in (data.get("text_contents") or []):
            text_content_map[txt_entry["placeholder_name"].lower()] = txt_entry

        # Sets for O(1) membership tests
        image_shape_names = set(image_shapes.keys())
        text_shape_names = set(text_shapes.keys())

        for shape_data in template_info['shapes']:
            shape_name = shape_data['name']

            if shape_name in image_shape_names:
                self._handle_image_shape(
                    slide, prs, shape_data,
                    image_source_map.get(shape_name.lower()),
                    embedded_images, result
                )
            elif shape_name in text_shape_names:
                self._handle_text_shape(
                    slide, shape_data,
                    text_content_map.get(shape_name.lower()),
                    result
                )
            else:
                self._recreate_shape(slide, shape_data)

    def _handle_image_shape(
        self,
        slide,
        prs: Presentation,
        shape_data: dict,
        img_entry: Optional[dict],
        embedded_images: Dict[str, ImageResult],
        result: SlideResult
    ) -> None:
        """Resolve and place a single image into a template shape."""
        img_result = None

        if img_entry:
            image_cell = img_entry.get("image_cell")
            image_source = img_entry.get("image_source")

            if image_cell and image_cell in embedded_images:
                img_result = embedded_images[image_cell]
            elif image_source and not image_source.startswith("http"):
                img_result = self.loader.load_from_path(image_source)

        if img_result and img_result.success:
            try:
                self._add_image_at_position(
                    slide, prs, img_result,
                    shape_data['left'], shape_data['top'],
                    shape_data['width'], shape_data['height']
                )
                result.has_image = True
            except Exception as e:
                result.image_error = str(e)
                self._add_placeholder_shape(slide, shape_data, "No Image")
        else:
            if img_result and img_result.error:
                result.image_error = img_result.error
            self._add_placeholder_shape(slide, shape_data, "No Image")

    def _handle_text_shape(
        self,
        slide,
        shape_data: dict,
        txt_entry: Optional[dict],
        result: SlideResult
    ) -> None:
        """Place text content into a template text shape."""
        text_content = txt_entry.get("text_content", []) if txt_entry else []
        try:
            self._add_text_from_template(slide, shape_data, text_content)
            result.text_added = True
        except Exception as e:
            err_msg = f"Text failed for '{shape_data['name']}': {e}"
            result.error = f"{result.error}; {err_msg}" if result.error else err_msg

    # ------------------------------------------------------------------
    # Legacy single-element slide population (pre-v8.0 backward compat)
    # ------------------------------------------------------------------

    def _populate_slide_legacy(
        self,
        slide,
        prs: Presentation,
        data: dict,
        embedded_images: Dict[str, ImageResult],
        template_info: dict,
        result: SlideResult
    ) -> None:
        """Populate a slide using legacy single image / single text fields."""
        text_content = data.get("text_content", [])
        image_source = data.get("image_source")
        image_cell = data.get("image_cell")

        # Resolve the single image
        img_result = None
        if image_cell and image_cell in embedded_images:
            img_result = embedded_images[image_cell]
        elif image_source and not image_source.startswith("http"):
            img_result = self.loader.load_from_path(image_source)

        for shape_data in template_info['shapes']:
            shape_name = shape_data['name']

            # Handle image placeholder
            if template_info['image_shape'] and shape_name == template_info['image_shape']['name']:
                if img_result and img_result.success:
                    try:
                        self._add_image_at_position(
                            slide, prs, img_result,
                            shape_data['left'], shape_data['top'],
                            shape_data['width'], shape_data['height']
                        )
                        result.has_image = True
                    except Exception as e:
                        result.image_error = str(e)
                        self._add_placeholder_shape(slide, shape_data, "No Image")
                else:
                    if img_result and img_result.error:
                        result.image_error = img_result.error
                    self._add_placeholder_shape(slide, shape_data, "No Image")

            # Handle text placeholder
            elif template_info['text_shape'] and shape_name == template_info['text_shape']['name']:
                try:
                    self._add_text_from_template(slide, shape_data, text_content)
                    result.text_added = True
                except Exception as e:
                    result.error = f"Text failed: {e}"

            # Copy other shapes as-is (backgrounds, decorations, etc.)
            else:
                self._recreate_shape(slide, shape_data)

    def _add_placeholder_shape(self, slide, shape_data: dict, text: str) -> None:
        """Add a placeholder rectangle shape."""
        if shape_data['type'] == MSO_SHAPE_TYPE.AUTO_SHAPE:
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                shape_data['left'], shape_data['top'],
                shape_data['width'], shape_data['height']
            )
            shape.name = shape_data['name']
            if shape.has_text_frame:
                shape.text_frame.paragraphs[0].text = text
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_image_at_position(
        self,
        slide,
        prs: Presentation,
        img_result: ImageResult,
        left, top, max_width, max_height
    ) -> None:
        """Add image at specified position, scaled to fit with alignment."""
        img_result.data.seek(0)
        orig_width, orig_height = self._get_image_dimensions(img_result.data)
        img_result.data.seek(0)

        # Calculate scaled size to fit within bounds
        final_width, final_height = self._calculate_scaled_size(
            orig_width, orig_height,
            max_width.inches, max_height.inches,
            IMG_SIZE_FIT_BOX
        )

        # Get alignment settings
        alignment = self.config.get_image_alignment()

        # Calculate position based on alignment within placeholder bounds
        img_left_inches, img_top_inches = self._calculate_image_position(
            final_width, final_height,
            left.inches, top.inches,
            max_width.inches, max_height.inches,
            alignment
        )

        slide.shapes.add_picture(
            img_result.data,
            Inches(img_left_inches), Inches(img_top_inches),
            Inches(final_width), Inches(final_height)
        )

    def _add_text_from_template(
        self,
        slide,
        shape_data: dict,
        text_content: List[Union[str, dict]]
    ) -> None:
        """Add text box following template formatting, populated with data."""
        textbox = slide.shapes.add_textbox(
            shape_data['left'], shape_data['top'],
            shape_data['width'], shape_data['height']
        )
        textbox.name = shape_data['name']

        tf = textbox.text_frame
        tf.word_wrap = True

        # Apply text overflow mode (v7.1 fix - was missing in template path)
        if self.config.text_overflow_mode == "shrink":
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        template_paragraphs = shape_data['paragraphs']

        # Build content mapping: template paragraph -> data
        # Template has: P0=ColC, P1=ColD, P2=empty, P3=ColE, P4=empty, P5=ColF
        # Data has: [{"column":"C","text":"..."}, {"column":"D","text":"..."}, ...]

        # Create a map of column letter to text
        content_map = {}
        for item in text_content:
            if isinstance(item, dict):
                col = item.get("column", "")
                text = item.get("text", "")
                content_map[col] = text
            else:
                # Legacy format - just use index
                pass

        # Dynamically map template paragraphs to user's columns (v7.0 fix)
        # Strategy: Detect spacer paragraphs (empty in template) and map user columns to content paragraphs

        # Build column sequence dynamically from user's text columns
        user_columns = list(content_map.keys())  # Columns that have data

        # Analyze template to find spacer positions (paragraphs with empty text in template)
        spacer_indices = set()
        for i, para_data in enumerate(template_paragraphs):
            # A paragraph is a spacer if all its runs have empty or whitespace-only text
            is_spacer = True
            for run_data in para_data.get('runs', []):
                if run_data.get('text', '').strip():
                    is_spacer = False
                    break
            if is_spacer and not para_data.get('runs'):
                is_spacer = True  # Empty paragraph with no runs is a spacer
            if is_spacer:
                spacer_indices.add(i)

        # Build dynamic column sequence: user columns go to non-spacer positions
        column_sequence = []
        user_col_idx = 0
        for i in range(len(template_paragraphs)):
            if i in spacer_indices:
                column_sequence.append("")  # Spacer
            elif user_col_idx < len(user_columns):
                column_sequence.append(user_columns[user_col_idx])
                user_col_idx += 1
            else:
                column_sequence.append("")  # No more user columns

        for i, para_data in enumerate(template_paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Determine what text to use
            if i < len(column_sequence):
                col_letter = column_sequence[i]
                if col_letter and col_letter in content_map:
                    new_text = content_map[col_letter]
                elif col_letter == "":
                    new_text = ""  # Keep empty paragraphs as spacers
                else:
                    new_text = ""  # No data for this column
            else:
                new_text = ""

            # Set alignment from template
            if para_data.get('alignment'):
                p.alignment = para_data['alignment']
            else:
                p.alignment = PP_ALIGN.CENTER

            # Add run with template formatting
            if para_data['runs']:
                run_data = para_data['runs'][0]  # Use first run's formatting
                run = p.add_run()
                run.text = new_text

                # Apply template formatting
                if run_data.get('font_name'):
                    run.font.name = run_data['font_name']
                if run_data.get('font_size'):
                    run.font.size = run_data['font_size']
                if run_data.get('bold') is not None:
                    run.font.bold = run_data['bold']
                if run_data.get('italic') is not None:
                    run.font.italic = run_data['italic']
                if run_data.get('color'):
                    run.font.color.rgb = run_data['color']
            else:
                # No template run formatting - just add text
                run = p.add_run()
                run.text = new_text

            # Apply configurable spacing
            p.space_after = Pt(self.config.paragraph_spacing)

    def _recreate_shape(self, slide, shape_data: dict) -> None:
        """Recreate a shape from template data (for backgrounds, etc.)."""
        # Only handle text boxes for now - skip complex shapes
        if shape_data['type'] == MSO_SHAPE_TYPE.TEXT_BOX:
            textbox = slide.shapes.add_textbox(
                shape_data['left'], shape_data['top'],
                shape_data['width'], shape_data['height']
            )
            textbox.name = shape_data['name']

            tf = textbox.text_frame
            tf.word_wrap = True

            for i, para_data in enumerate(shape_data['paragraphs']):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                if para_data.get('alignment'):
                    p.alignment = para_data['alignment']

                for run_data in para_data['runs']:
                    run = p.add_run()
                    run.text = run_data['text']
                    if run_data.get('font_name'):
                        run.font.name = run_data['font_name']
                    if run_data.get('font_size'):
                        run.font.size = run_data['font_size']
                    if run_data.get('bold') is not None:
                        run.font.bold = run_data['bold']

    def _create_slide(
        self,
        prs: Presentation,
        data: dict,
        embedded_images: Dict[str, ImageResult]
    ) -> SlideResult:
        """Create a single slide (original blank slide method)."""
        index = data.get("row_index", len(prs.slides))
        result = SlideResult(index=index, success=True)

        try:
            # Add blank slide (layout 6 is typically blank)
            slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
            slide = prs.slides.add_slide(slide_layout)

            # Handle image
            image_source = data.get("image_source")
            image_cell = data.get("image_cell")

            img_result = None

            if image_cell and image_cell in embedded_images:
                img_result = embedded_images[image_cell]
            elif image_source and not image_source.startswith("http"):
                img_result = self.loader.load_from_path(image_source)

            if img_result:
                if img_result.success:
                    try:
                        self._add_image(slide, prs, img_result)
                        result.has_image = True
                    except Exception as e:
                        result.image_error = str(e)
                        logger.warning(f"Slide {index}: Image add failed: {e}")
                else:
                    result.image_error = img_result.error
                    logger.warning(f"Slide {index}: Image load failed: {img_result.error}")

            # Add text
            text_content = data.get("text_content", [])
            if text_content:
                try:
                    self._add_text(slide, prs, text_content)
                    result.text_added = True
                except Exception as e:
                    result.error = f"Text add failed: {e}"
                    logger.warning(f"Slide {index}: Text add failed: {e}")

        except Exception as e:
            result.success = False
            result.error = str(e)
            logger.error(f"Slide {index}: Creation failed: {e}", exc_info=True)

        return result

    def _get_image_dimensions(self, img_data: BytesIO) -> Tuple[int, int]:
        """Get image dimensions in pixels."""
        img_data.seek(0)
        with Image.open(img_data) as img:
            return img.size

    def _calculate_scaled_size(
        self,
        orig_width: int,
        orig_height: int,
        max_width: float,
        max_height: float,
        mode: str
    ) -> Tuple[float, float]:
        """Calculate scaled image dimensions based on sizing mode."""
        aspect_ratio = orig_width / orig_height if orig_height > 0 else 1.0

        if mode == IMG_SIZE_STRETCH:
            return max_width, max_height
        elif mode == IMG_SIZE_FIT_WIDTH:
            return max_width, max_width / aspect_ratio
        elif mode == IMG_SIZE_FIT_HEIGHT:
            return max_height * aspect_ratio, max_height
        else:  # IMG_SIZE_FIT_BOX
            width_if_fit_height = max_height * aspect_ratio
            height_if_fit_width = max_width / aspect_ratio

            if width_if_fit_height <= max_width:
                return width_if_fit_height, max_height
            else:
                return max_width, height_if_fit_width

    def _calculate_image_position(
        self,
        img_width: float,
        img_height: float,
        box_left: float,
        box_top: float,
        box_width: float,
        box_height: float,
        alignment: ImageAlignment
    ) -> Tuple[float, float]:
        """
        Calculate image position based on alignment within bounding box.

        Args:
            img_width: Scaled image width in inches
            img_height: Scaled image height in inches
            box_left: Bounding box left in inches
            box_top: Bounding box top in inches
            box_width: Bounding box width in inches
            box_height: Bounding box height in inches
            alignment: ImageAlignment with vertical and horizontal settings

        Returns:
            Tuple of (left_inches, top_inches)
        """
        # Calculate horizontal position
        if alignment.horizontal == IMG_ALIGN_LEFT:
            left = box_left
        elif alignment.horizontal == IMG_ALIGN_RIGHT:
            left = box_left + box_width - img_width
        else:  # center (default)
            left = box_left + (box_width - img_width) / 2

        # Calculate vertical position
        if alignment.vertical == IMG_ALIGN_TOP:
            top = box_top
        elif alignment.vertical == IMG_ALIGN_BOTTOM:
            top = box_top + box_height - img_height
        else:  # center (default)
            top = box_top + (box_height - img_height) / 2

        return (left, top)

    def _add_image(
        self,
        slide,
        prs: Presentation,
        img_result: ImageResult
    ) -> None:
        """Add image to slide with configurable sizing and alignment."""
        img_result.data.seek(0)
        orig_width, orig_height = self._get_image_dimensions(img_result.data)
        img_result.data.seek(0)

        final_width, final_height = self._calculate_scaled_size(
            orig_width, orig_height,
            self.config.img_width, self.config.img_height,
            self.config.img_size_mode
        )

        # Get alignment settings
        alignment = self.config.get_image_alignment()

        # Calculate bounding box based on horizontal alignment
        slide_width_inches = prs.slide_width.inches
        if alignment.horizontal == IMG_ALIGN_LEFT:
            box_left = self.config.img_left
        elif alignment.horizontal == IMG_ALIGN_RIGHT:
            box_left = max(0, slide_width_inches - self.config.img_width - self.config.img_left)
        else:  # center (legacy default)
            box_left = (slide_width_inches - self.config.img_width) / 2
        box_top = self.config.img_top
        box_width = self.config.img_width
        box_height = self.config.img_height

        # Calculate position based on alignment
        img_left, img_top = self._calculate_image_position(
            final_width, final_height,
            box_left, box_top, box_width, box_height,
            alignment
        )

        pic = slide.shapes.add_picture(
            img_result.data,
            Inches(img_left),
            Inches(img_top),
            width=Inches(final_width),
            height=Inches(final_height)
        )

    def _add_text(
        self,
        slide,
        prs: Presentation,
        text_content: List[Union[str, dict]]
    ) -> None:
        """Add text content to slide with per-column formatting and configurable spacing.

        Supports two modes:
        - Auto flow: All columns in single textbox (default/legacy)
        - Fixed positions: Columns with fixed positions get separate textboxes
        """
        # Separate auto and fixed columns
        auto_items = []
        fixed_items = []

        for item in text_content:
            if isinstance(item, dict):
                col = item.get("column", "")
                col_pos = self.config.get_column_position(col)
                if col_pos and col_pos.mode == "fixed" and col_pos.top is not None:
                    fixed_items.append((item, col_pos))
                else:
                    auto_items.append(item)
            else:
                auto_items.append(item)

        # Add auto-flow columns in single textbox (original behavior)
        if auto_items:
            self._add_text_auto_flow(slide, prs, auto_items)

        # Add fixed-position columns as separate textboxes
        for item, col_pos in fixed_items:
            self._add_text_fixed(slide, prs, item, col_pos)

    def _add_text_auto_flow(
        self,
        slide,
        prs: Presentation,
        text_items: List[Union[str, dict]]
    ) -> None:
        """Add text items in auto-flow mode (single textbox, sequential paragraphs)."""
        text_height = prs.slide_height.inches - self.config.text_top - 0.5
        text_width = max(1.0, prs.slide_width.inches - self.config.text_left - 0.5)

        textbox = slide.shapes.add_textbox(
            Inches(self.config.text_left),
            Inches(self.config.text_top),
            Inches(text_width),
            Inches(text_height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Set text overflow mode (NEW in v6.2)
        if self.config.text_overflow_mode == "shrink":
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for i, item in enumerate(text_items):
            if isinstance(item, dict):
                text = item.get("text", "")
                col_format = self.config.get_column_format(item.get("column", ""))
            else:
                text = item
                col_format = ColumnFormat(column="", font_size=self.config.font_size)

            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            run = p.add_run()
            run.text = text

            font = run.font
            font.size = Pt(col_format.font_size)
            font.bold = col_format.bold
            font.italic = col_format.italic
            font.name = col_format.font_name
            font.color.rgb = col_format.get_rgb_color()

            p.alignment = self.config.get_text_pp_align()
            p.space_after = Pt(self.config.paragraph_spacing)

    def _add_text_fixed(
        self,
        slide,
        prs: Presentation,
        item: dict,
        col_pos: ColumnPosition
    ) -> None:
        """Add a single text item at a fixed position."""
        text = item.get("text", "")
        col = item.get("column", "")
        col_format = self.config.get_column_format(col)

        # Calculate dimensions
        left = col_pos.left
        top = col_pos.top
        width = col_pos.width if col_pos.width else (prs.slide_width.inches - 1.0)
        height = 1.5  # Default height for fixed text boxes

        textbox = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Set text overflow mode (NEW in v6.2)
        if self.config.text_overflow_mode == "shrink":
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        font = run.font
        font.size = Pt(col_format.font_size)
        font.bold = col_format.bold
        font.italic = col_format.italic
        font.name = col_format.font_name
        font.color.rgb = col_format.get_rgb_color()

        p.alignment = self.config.get_text_pp_align()
        p.space_after = Pt(self.config.paragraph_spacing)


def create_presentation(
    slide_data: List[dict],
    config: Optional[SlideConfig] = None,
    embedded_images: Optional[Dict[str, ImageResult]] = None,
    template_file: Optional[Union[bytes, BytesIO, str, Path]] = None,
    progress_callback: Optional[Callable[[str, int, int], None]] = None
) -> GenerationResult:
    """Convenience function to create a presentation."""
    generator = PPTXGenerator(config)
    return generator.generate(slide_data, embedded_images, template_file, progress_callback)
